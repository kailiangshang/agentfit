#!/usr/bin/env python3
"""Deterministic compiler from the AgentFit HTML slide sources to a native PPTX.

Every slide is authored as absolutely positioned elements over a 1280x720
canvas (see slides/common.css). This compiler maps each element to a native,
editable PowerPoint shape at identical coordinates and colors:

- .panel / .pill / .stamp / .chip  -> rounded-rectangle autoshape
- .top-rule / .line / .line-dim / .line-v -> thin filled rectangle
- .node                             -> oval
- text elements (p, or div with text) -> text box with font metrics taken
  from inline styles and the common.css class table below

1 CSS pixel maps to exactly 9525 EMU (1280px -> 13.333in, 720px -> 7.5in),
so geometry is preserved without scaling.
"""

from __future__ import annotations

import re
import sys
from html.parser import HTMLParser
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parent
SLIDES_DIR = ROOT / "slides"
OUTPUT_PATH = ROOT / "agentfit-submission.pptx"
EXPECTED_SLIDES = 17

PX_TO_EMU = 9525
PAGE_W_EMU = 1280 * PX_TO_EMU
PAGE_H_EMU = 720 * PX_TO_EMU

# class -> (font-size px, color, bold, mono) defaults from slides/common.css
CLASS_TEXT_DEFAULTS = {
    "eyebrow": (14, "#1a8d85", True, True),
    "title": (42, "#102a43", True, False),
    "subtitle": (20, "#53697c", False, False),
    "label": (14, "#53697c", True, True),
    "body-copy": (18, "#53697c", False, False),
    "small-copy": (14, "#53697c", False, False),
    "tiny-copy": (12, "#718190", False, False),
    "section-title": (24, "#102a43", True, False),
    "big-number": (52, "#1a8d85", True, True),
    "code-tag": (13, "#1a8d85", True, True),
    "footer-left": (13, "#718190", False, True),
    "footer-right": (13, "#718190", False, True),
}

TEXT_CLASSES = set(CLASS_TEXT_DEFAULTS)
DECOR_CLASSES = {
    "panel", "top-rule", "line", "line-v", "line-dim", "pill", "node",
    "stamp", "chip",
}

HEX_RE = re.compile(r"#([0-9a-fA-F]{3}|[0-9a-fA-F]{6})\b")
STYLE_RE = re.compile(r"([a-z-]+)\s*:\s*([^;]+)")
CSS_RULE_RE = re.compile(r"([^{}]+)\{([^{}]*)\}")


def load_css(path: Path) -> dict[str, dict[str, str]]:
    """Map each class name (and 'parent child' combos) to its declarations."""
    css: dict[str, dict[str, str]] = {}
    for selector, body in CSS_RULE_RE.findall(path.read_text(encoding="utf-8")):
        for name in selector.strip().split(","):
            classes = re.findall(r"\.([A-Za-z-]+)", name)
            key = " ".join(classes)
            css.setdefault(key, {}).update(parse_style(body))
    return css


def parse_style(style: str) -> dict[str, str]:
    return {k.strip(): v.strip() for k, v in STYLE_RE.findall(style or "")}


def to_rgb(value: str) -> RGBColor | None:
    match = HEX_RE.search(value or "")
    if not match:
        return None
    hexpart = match.group(1)
    if len(hexpart) == 3:
        hexpart = "".join(ch * 2 for ch in hexpart)
    return RGBColor(int(hexpart[0:2], 16), int(hexpart[2:4], 16), int(hexpart[4:6], 16))


def length(styles: dict[str, str], name: str, default: float | None = None) -> float | None:
    raw = styles.get(name)
    if raw is None:
        return default
    match = re.match(r"(-?[\d.]+)px", raw.strip())
    return float(match.group(1)) if match else default


class Node:
    __slots__ = ("tag", "attrs", "content")

    def __init__(self, tag: str, attrs: dict[str, str]):
        self.tag = tag
        self.attrs = attrs
        # ordered inline content: ("#text", str) or child Node
        self.content: list = []

    @property
    def children(self) -> list["Node"]:
        return [item for item in self.content if not isinstance(item, tuple)]

    @property
    def classes(self) -> set[str]:
        return set((self.attrs.get("class") or "").split())

    @property
    def style(self) -> dict[str, str]:
        return parse_style(self.attrs.get("style", ""))

    def direct_text(self) -> str:
        return "".join(item[1] for item in self.content if isinstance(item, tuple))


class SlideParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.root = Node("root", {})
        self.stack: list[Node] = [self.root]

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        node = Node(tag, {k: (v or "") for k, v in attrs})
        self.stack[-1].content.append(node)
        if tag not in ("br", "img", "meta", "link"):
            self.stack.append(node)

    def handle_startendtag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        self.stack[-1].content.append(Node(tag, {k: (v or "") for k, v in attrs}))

    def handle_endtag(self, tag: str) -> None:
        for index in range(len(self.stack) - 1, 0, -1):
            if self.stack[index].tag == tag:
                del self.stack[index:]
                break

    def handle_data(self, data: str) -> None:
        self.stack[-1].content.append(("#text", data))


# ---------------------------------------------------------------- rendering

def is_positioned(node: Node, dark: bool) -> bool:
    classes = node.classes
    if classes & TEXT_CLASSES or classes & DECOR_CLASSES:
        return True
    styles = node.style
    return "left" in styles or "top" in styles or "position" in styles


def add_shape(slide, kind: str, box: dict, styles: dict, classes: set, dark: bool) -> None:
    left = Emu(int(box["left"] * PX_TO_EMU))
    top = Emu(int(box["top"] * PX_TO_EMU))
    width = Emu(int(box["width"] * PX_TO_EMU))
    height = Emu(int(box["height"] * PX_TO_EMU))

    if "panel" in classes:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    elif "node" in classes:
        shape_type = MSO_SHAPE.OVAL
    elif {"pill", "stamp", "chip"} & classes:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:  # rules and lines
        shape_type = MSO_SHAPE.RECTANGLE

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False

    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            short = min(box["width"], box["height"])
            if short > 0:
                shape.adjustments[0] = min(0.5, 18.0 / short)
        except Exception:
            pass

    background = to_rgb(styles.get("background", ""))
    border = to_rgb(styles.get("border-color", ""))
    if background is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = background
    elif "panel-dark" in classes:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x13, 0x2F, 0x47)
    elif "line" in classes or "top-rule" in classes:
        color = border or to_rgb("#1a8d85")
        shape.fill.solid()
        shape.fill.fore_color.rgb = color or RGBColor(0x1A, 0x8D, 0x85)
    else:
        color_map = {
            "soft-teal": "#dceeea", "soft-coral": "#f8ddd4",
            "soft-amber": "#f3e5b8", "soft-blue": "#dfeaf2",
            "soft-ink": "#eef1f5",
        }
        fill_hex = next((v for k, v in color_map.items() if k in classes), None)
        if fill_hex:
            shape.fill.solid()
            shape.fill.fore_color.rgb = to_rgb(fill_hex)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    border_color = border or to_rgb(styles.get("background", ""))
    line = shape.line
    if {"top-rule", "line", "line-v", "line-dim"} & classes:
        line.fill.background()
    else:
        if border_color is not None:
            line.color.rgb = border_color
            line.width = Pt(1.5)
        else:
            line.fill.background()


def collect_tokens(node: Node, inherited_bold: bool = False) -> list:
    """Flatten inline content to (text, bold) tokens; None marks a line break."""
    self_bold = inherited_bold or node.tag in ("b", "strong")
    result: list = []
    for item in node.content:
        if isinstance(item, tuple):
            text = re.sub(r"\s+", " ", item[1])
            if text:
                result.append((text, self_bold))
        elif item.tag == "br":
            result.append(None)
        else:
            result.extend(collect_tokens(item, self_bold))
    return result


def add_text(slide, node: Node, box: dict, styles: dict, classes: set, dark: bool) -> None:
    tokens = collect_tokens(node)
    if not any(isinstance(t, tuple) and t[0].strip() for t in tokens):
        return

    left = Emu(int(box["left"] * PX_TO_EMU))
    top = Emu(int(box["top"] * PX_TO_EMU))
    width = Emu(int(max(box.get("width") or 400, 16) * PX_TO_EMU))
    height = Emu(int(max(box.get("height") or 20, 14) * PX_TO_EMU))

    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    text_class = next((c for c in classes if c in CLASS_TEXT_DEFAULTS), None)
    default_size, default_color, default_bold, mono = CLASS_TEXT_DEFAULTS.get(
        text_class, (16, "#53697c", False, False)
    )
    if dark:
        if text_class == "title":
            default_color = "#ffffff"
        elif text_class == "subtitle":
            default_color = "#c6d3dd"
        elif default_color == "#102a43":
            default_color = "#ffffff"

    font_size = styles.get("font-size", "")
    match = re.match(r"([\d.]+)px", font_size)
    size = float(match.group(1)) if match else default_size
    color = to_rgb(styles.get("color", "")) or to_rgb(default_color)
    bold = default_bold

    centered = "text-align:center" in (node.attrs.get("style") or "").replace(" ", "")

    lines: list[list[tuple[str, bool]]] = [[]]
    for token in tokens:
        if token is None:
            lines.append([])
            continue
        lines[-1].append(token)

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        if centered:
            paragraph.alignment = PP_ALIGN.CENTER
        if not line:
            run = paragraph.add_run()
            run.text = " "
            run.font.size = Pt(size)
            continue
        for segment, segment_bold in line:
            run = paragraph.add_run()
            run.text = segment
            run.font.size = Pt(size)
            if color is not None:
                run.font.color.rgb = color
            run.font.bold = bold or segment_bold
            run.font.name = (
                "Courier New" if mono or "mono" in classes else "Noto Sans CJK SC"
            )


CSS_TABLE: dict[str, dict[str, str]] = {}


def effective_style(node: Node, classes: list[str], dark: bool, appendix: bool) -> dict[str, str]:
    merged: dict[str, str] = {}
    ordered = []
    if dark or appendix:
        prefix = ("appendix " if appendix else "") + ("dark " if dark else "")
        # 组合类 (.appendix .title / .dark .subtitle) 优先于基础类
        for cls in classes:
            combo = f"{prefix}{cls}"
            if combo in CSS_TABLE:
                ordered.append(combo)
    ordered += classes
    for key in ordered:
        if key in CSS_TABLE:
            merged.update(CSS_TABLE[key])
    merged.update(node.style)
    return merged


def visual_width(text: str, font_size: float) -> float:
    units = sum(1.0 if ord(ch) > 0x2E80 else 0.62 for ch in text)
    return units * font_size


def box_of(node: Node, styles: dict[str, str], font_size: float, line_count: int,
           est_width: float | None = None) -> dict:
    left = length(styles, "left")
    top = length(styles, "top")
    width = length(styles, "width")
    height = length(styles, "height")
    right = length(styles, "right")
    bottom_px = length(styles, "bottom")
    if left is None and right is not None and width is not None:
        left = 1280 - right - width
    if top is None and bottom_px is not None and height is not None:
        top = 720 - bottom_px - height
    if width is None:
        if est_width:
            width = min(1280 - (left or 0), max(16, est_width + 8))
        else:
            width = max(40, 1280 - (left or 0))
    if height is None:
        line_height = 1.3
        lh_raw = styles.get("line-height", "")
        lh_match = re.match(r"([\d.]+)", lh_raw)
        if lh_match:
            line_height = float(lh_match.group(1))
        height = max(14, line_count * font_size * line_height)
    return {"left": left or 0, "top": top or 0, "width": width, "height": height}


def render_element(slide, node: Node, dark: bool, appendix: bool, parent_box: dict | None) -> None:
    class_list = [c for c in (node.attrs.get("class") or "").split()]
    classes = set(class_list)
    eff = effective_style(node, class_list, dark, appendix)
    has_own_position = "left" in eff or "top" in eff or "right" in eff or "bottom" in eff
    nested = parent_box is not None and not has_own_position and bool(classes & TEXT_CLASSES or node.tag == "p")

    tokens = collect_tokens(node)
    text_line_count = sum(1 for token in tokens if token is None) + 1 if tokens else 1
    font_size_px = float(re.match(r"([\d.]+)", eff.get("font-size", "16")).group(1))
    est_w = None
    if tokens:
        line_texts = [""]
        for token in tokens:
            if token is None:
                line_texts.append("")
            else:
                line_texts[-1] += token[0]
        est_w = max(visual_width(txt, font_size_px) for txt in line_texts)
    box = box_of(node, eff, font_size_px, text_line_count, est_w)

    if nested and parent_box is not None:
        margin_top = length(eff, "margin-top", 4) or 4
        margin_side = length(eff, "margin-left", 8) or 8
        box["left"] = parent_box["left"] + margin_side
        box["top"] = parent_box["top"] + margin_top
        box["width"] = parent_box["width"] - 2 * margin_side

    is_decor = bool(classes & DECOR_CLASSES)
    has_own_text = bool(classes & TEXT_CLASSES) or node.tag == "p" or bool(tokens)

    if is_decor and tokens:
        # 装饰容器（pill/chip/stamp/node）：底板形状 + 文本框，字体继承内嵌段落
        add_shape(slide, node.tag, box, eff, classes, dark)
        first_text_child = next(
            (c for c in node.children if c.tag == "p" and collect_tokens(c)), None
        )
        text_styles = eff
        if first_text_child is not None:
            child_eff = effective_style(first_text_child, [], dark, appendix)
            merged = dict(eff)
            merged.update(first_text_child.style)
            text_styles = merged
        margin_top = length(text_styles, "margin-top", 4) or 4
        margin_side = length(text_styles, "margin-left", 8) or 8
        fs_match = re.match(r"([\d.]+)", text_styles.get("font-size", "16"))
        decor_font = float(fs_match.group(1)) if fs_match else 16.0
        decor_lines = sum(1 for token in tokens if token is None) + 1
        lh_match = re.match(r"([\d.]+)", text_styles.get("line-height", "1.3"))
        decor_lh = float(lh_match.group(1)) if lh_match else 1.3
        text_box = {
            "left": box["left"] + margin_side,
            "top": box["top"] + margin_top,
            "width": box["width"] - 2 * margin_side,
            "height": max(14, decor_lines * decor_font * decor_lh),
        }
        add_text(slide, node, text_box, text_styles, classes, dark)
        return
    if is_decor and not tokens:
        add_shape(slide, node.tag, box, eff, classes, dark)
        child_parent = box
    elif has_own_text:
        # 内联内容(含嵌套段落文本)已并入 tokens，不再递归，避免重复文本框
        add_text(slide, node, box, eff, classes, dark)
        return
    else:
        add_shape(slide, node.tag, box, eff, classes, dark)
        child_parent = box

    for child in node.children:
        if child.tag == "br":
            continue
        render_element(slide, child, dark, appendix, child_parent)


def find_body(root: Node) -> Node:
    stack = [root]
    while stack:
        node = stack.pop()
        if node.tag == "body":
            return node
        stack.extend(node.children)
    return root


def compile_slide(presentation: Presentation, html_path: Path) -> None:
    source = html_path.read_text(encoding="utf-8")
    parser = SlideParser()
    parser.feed(source)
    body = find_body(parser.root)
    dark = "dark" in body.classes
    appendix = "appendix" in body.classes

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = (
        RGBColor(0x0B, 0x22, 0x36) if dark else RGBColor(0xF6, 0xF2, 0xE8)
    )

    for child in body.children:
        if child.tag == "br":
            continue
        render_element(slide, child, dark, appendix, None)


def build(output_path: Path = OUTPUT_PATH) -> None:
    global CSS_TABLE
    CSS_TABLE = load_css(SLIDES_DIR / "common.css")
    slide_files = sorted(SLIDES_DIR.glob("[0-9][0-9]-*.html"))
    if len(slide_files) != EXPECTED_SLIDES:
        raise RuntimeError(
            f"expected {EXPECTED_SLIDES} HTML slides, found {len(slide_files)}"
        )
    presentation = Presentation()
    presentation.slide_width = Emu(PAGE_W_EMU)
    presentation.slide_height = Emu(PAGE_H_EMU)
    for html_path in slide_files:
        compile_slide(presentation, html_path)
    presentation.save(output_path)


def main() -> int:
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else OUTPUT_PATH
    build(target)
    print(f"compiled {EXPECTED_SLIDES} slides -> {target}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
