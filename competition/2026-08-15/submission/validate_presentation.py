#!/usr/bin/env python3
"""Validate the final AgentFit preliminary-submission deliverables."""

from __future__ import annotations

import argparse
import re
import sys
import zipfile
from xml.etree import ElementTree
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader


EXPECTED_SLIDES = 17

EXPECTED_PAGE_MARKERS = tuple(
    f"{index:02d} / 12" for index in range(1, 13)
) + tuple(f"A{index} / 05" for index in range(1, 6))

EXPECTED_META_AGENT_IDENTITIES = (
    "EngagementLead",
    "BusinessEngineer",
    "AgentArchitect",
    "ValidationEngineer",
    "GovernanceAuditor",
)

EXPECTED_SKILL_ENTRIES = (
    "任务编译",
    "能力对齐",
    "候选建图",
    "统一试验",
    "独立审计",
    "人工门禁",
    "经验沉淀",
)

SLIDE_11_EVIDENCE_STATEMENT = (
    "OpsPilot 是官方案例锚点；retail / airline 仅为探索性 Demo，"
    "使用非官方 evaluator，不是正式 Candidate，也不是官方分数。"
)

EVALUATION_IDENTITY = "CandidateVersion × SampleVersion × RunIndex"

EXPECTED_PAGE_TITLES = (
    "企业真正缺少的，是选对 Agent 方案。",
    "OpsPilot 官方示例：4 个 Worker 加 1 个 Leader，仍未回答该用哪种。",
    "AgentFit 把 baseline 材料编译成首个 ProjectCase。",
    "两个事故变成 TaskSample：输入、输出与验收。",
    "方案空间按四层资产组织，调整逐层受控。",
    "样本驱动的持续学习：离散信用分配，回归防遗忘。",
    "五个元 Agent 组成 AgentFit Learning Loop，分别负责目标、样本、方案、实验与诊断。",
    "AgentTeams 承载 Worker、Team、Room、Human；AgentFit 落地 Dossier 与 Trace。",
    "追溯四件套：实体分组、hash 链账本、检查器族、泳道报告。",
    "交付的不是一张架构图，而是可复现的 AgentSolutionPackage。",
    "证据账本：OpsPilot 为官方锚点，retail/airline 仅探索性 Demo。",
    "从 OpsPilot 回到通用：Fit 是有证据选对方案。",
    "七层 ML 映射升级为硬映射：持续学习视角。",
    "五个 Agent Identity：判断权、状态边界与责任产物。",
    "七个 Skill、Tool 与 MCP/HTTP 契约、Memory 与 Trace。",
    "Human 门禁、风险、异常与回滚。",
    "开放、依赖、许可证、baseline 引用与未实现边界。",
)

SLIDE_REQUIRED_TERMS = {
    2: ("OpsPilot", "ProjectCase", "baseline"),
    4: ("TaskSample", "db_pool_exhausted", "slow_sql_degradation"),
    5: (
        "完整方案空间",
        "四层资产",
        "L1 · SOLID",
        "source tool",
        "触达纪律",
        "复杂度控制",
        "工具",
        "Skill",
        "MCP",
        "Memory",
        "模型",
        "Agent 拓扑",
        "Human 边界",
        "C0",
        "C1",
        "C2",
        "C3",
    ),
    6: (
        "样本驱动的持续学习",
        "离散信用分配",
        "样本分类",
        "层内更新",
        "回归验证",
        "资产沉淀",
        "漂移",
        "五阶段闭环",
        "机器学习的工程纪律",
        "样本工程",
        "最小方案假设",
        "批量试验",
        "Trace 误差分析",
        "新案例验证 / 停止",
        "RegressionPool",
    ),
    7: (
        "AgentFit Learning Loop",
        "目标 / 停止控制",
        "业务理解 / 样本工程",
        "方案建模 / 结构选择",
        "批量试验 / 指标采集",
        "误差分析 / 治理审计",
        "EngagementLead",
        "BusinessEngineer",
        "AgentArchitect",
        "ValidationEngineer",
        "GovernanceAuditor",
    ),
    8: ("Worker", "Team", "Room", "Human", "Dossier", "Trace"),
    9: ("实体分组", "HASH 链账本", "检查器族", "泳道报告", "Skill", "RegressionPool"),
    10: (
        "可复现的 AgentSolutionPackage",
        "方案版本",
        "实验历史",
        "失败记录",
    ),
    11: (
        SLIDE_11_EVIDENCE_STATEMENT,
    ),
    13: (
        "七层 ML 映射",
        "G, Π, θ, ρ",
        "层内拟合",
        "候选对照",
        "Meta-learning",
    ),
    14: (
        "EngagementLead",
        "BusinessEngineer",
        "AgentArchitect",
        "ValidationEngineer",
        "GovernanceAuditor",
    ),
}

FIRST_PAGE_ML_BAN = (
    "Architecture Search",
    "NAS",
    "Meta-learning",
    "Meta learning",
    "七层",
    "L1",
    "L2",
    "L3",
    "SVD",
    "argmin",
    "Pareto",
    "贝叶斯",
    "inner loop",
    "outer loop",
)

REQUIRED_TERMS = (
    "AgentFit",
    "AgentTeams",
    "OpsPilot",
    "ProjectCase",
    "db_pool_exhausted",
    "slow_sql_degradation",
    "TaskSample",
    "Episode",
    EVALUATION_IDENTITY,
    "EngagementLead",
    "BusinessEngineer",
    "AgentArchitect",
    "ValidationEngineer",
    "GovernanceAuditor",
    "Agentless",
    "C0",
    "C1",
    "C2",
    "C3",
    "Human",
    "Skill",
    "MCP",
    "HTTP",
    "上下文",
    "验证",
    "安全",
    "开放",
    "未实现",
    "AgentSolutionPackage",
    "设计契约，非运行证据",
    "requires_runtime_trial",
    "Worker",
    "Team",
    "Room",
    "Dossier",
    "同一搜索空间",
    "七层 ML 映射",
    "元 Agent",
    "业务执行 Agent",
)

BROADENED_FORBIDDEN_TERMS = (
    "AutoML",
    "automl",
    "Auto ML",
    "自动机器学习",
    "semantic gradient",
    "semantic gradients",
    "semantic backpropagation",
    "语义梯度",
    "backpropagation",
    "back propagation",
    "back-propagation",
    "反向传播",
    "official benchmark accuracy",
    "official benchmark score",
    "official benchmark results",
    "官方 benchmark 准确率",
    "官方 benchmark accuracy",
    "官方 benchmark 分数",
    "官方 benchmark score",
    "官方 benchmark 结果",
    "官方基准准确率",
    "官方评测准确率",
    "正式 Candidate 已完成",
    "正式 Candidate 执行完成",
    "正式 Candidate 运行完成",
    "正式 Candidate 已跑通",
    "Candidate execution completed",
    "formal Candidate completed",
    "formal Candidate execution completed",
    "正式候选已完成",
    "正式候选执行完成",
    "正式候选运行完成",
    "正式候选已跑通",
    "AgentTeams 端到端集成完成",
    "AgentTeams 端到端集成已完成",
    "AgentTeams 已端到端集成",
    "AgentTeams 已完成端到端集成",
    "AgentTeams end-to-end integration complete",
    "AgentTeams end-to-end integration completed",
    "AgentTeams end-to-end integrated",
    "AgentTeams 跑通最小闭环",
)

FORBIDDEN_DECLARATION_PATTERNS = (
    (
        "verified Meta-learning claim",
        re.compile(
            r"\bmeta[- ]?learning\b\s+"
            r"(?:(?:is|was|has\s+been|have\s+been)\s+)?"
            r"(?:verified|validated|confirmed|proven)\b"
            r"|\b(?:verified|validated|confirmed|proven)\b\s+"
            r"\bmeta[- ]?learning\b"
            r"|(?:已验证|验证完成|已证实|验证过)"
            r"(?:的)?\s*meta[- ]?learning"
            r"|meta[- ]?learning\s*(?:已验证|验证完成|已证实|验证过)",
        ),
    ),
)

FORBIDDEN_TERMS = (
    "Agent 方案训练系统",
    "AutoML for Agents",
    "语义反向传播",
    "已选定 C2",
    "已选定 C1",
    "C2 胜出",
    "已跑通最小闭环",
    "ROI 提升",
    "准确率 9",
    "ImageNet",
    "90%+",
    "Meta-learning 已验证",
    "已开放",
    "proxy 分数是官方结果",
    "代理分数是官方结果",
    "proxy score is official",
    "exploratory proxy scores are official results",
) + BROADENED_FORBIDDEN_TERMS

PRESENTATIONML_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"


def forbidden_declarations(text: str) -> list[str]:
    """Return normalized literal and regex forbidden declarations in text."""
    normalized = _normalized(text).casefold()
    claim_text = re.sub(r"\s+", " ", text).casefold()
    declarations = [
        term
        for term in FORBIDDEN_TERMS
        if _normalized(term).casefold() in normalized
    ]
    declarations.extend(
        label
        for label, pattern in FORBIDDEN_DECLARATION_PATTERNS
        if pattern.search(claim_text)
    )
    return list(dict.fromkeys(declarations))


def _shape_text(shape: object) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        texts.append(shape.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                texts.append(cell.text)
    return texts


def _iter_shapes(shapes: object):
    for shape in shapes:
        yield shape
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_shapes(nested_shapes)


def _slide_text(slide: object) -> str:
    return "\n".join(
        text
        for shape in _iter_shapes(slide.shapes)
        for text in _shape_text(shape)
        if text.strip()
    )


def _normalized(text: str) -> str:
    return re.sub(r"\s+", "", text)


def _required_term_present(text: str, term: str) -> bool:
    """Match exact terms, allowing only the evaluation identity to span lines."""
    if term in text:
        return True
    return term == EVALUATION_IDENTITY and _normalized(term) in _normalized(text)


def _has_visible_native_content(
    slide: object, slide_width: int, slide_height: int
) -> bool:
    for shape in _iter_shapes(slide.shapes):
        if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA):
            continue
        left = getattr(shape, "left", 0)
        top = getattr(shape, "top", 0)
        width = getattr(shape, "width", 0)
        height = getattr(shape, "height", 0)
        if width <= 0 or height <= 0:
            continue
        if left >= slide_width or top >= slide_height:
            continue
        if left + width <= 0 or top + height <= 0:
            continue
        if any(_normalized(text) for text in _shape_text(shape)):
            return True
    return False


def _package_editability_errors(pptx_path: Path) -> list[str]:
    errors: list[str] = []
    try:
        with zipfile.ZipFile(pptx_path) as package:
            media = sorted(
                name
                for name in package.namelist()
                if name.startswith("ppt/media/") and not name.endswith("/")
            )
            for name in media:
                errors.append(f"PPTX package contains embedded media: {name}")
            for slide_number in range(1, EXPECTED_SLIDES + 1):
                member = f"ppt/slides/slide{slide_number}.xml"
                if member not in package.namelist():
                    continue
                root = ElementTree.fromstring(package.read(member))
                transition = root.find(
                    f"{{{PRESENTATIONML_NAMESPACE}}}transition"
                )
                if transition is not None:
                    errors.append(f"PPTX slide {slide_number} contains a transition")
    except (OSError, zipfile.BadZipFile, ElementTree.ParseError) as exc:
        errors.append(f"Unable to inspect PPTX package {pptx_path}: {exc}")
    return errors


NUMBERED_ENTRY_PATTERN = re.compile(
    r"(?<![\w/])(\d+)(?:\s*/\s*|\s+)(?!\d+\b)([^\s/·•]+)"
)


def _structured_integer_entries(text: str) -> list[tuple[int, str]]:
    """Extract numbered entries with direct or slash separators, not page markers."""
    return [
        (int(number), label)
        for number, label in NUMBERED_ENTRY_PATTERN.findall(text)
        if label not in {"/", "·", "•"}
    ]


def _skill_list_residual(text: str) -> str:
    """Return non-entry content from an already isolated Skill list."""
    text = re.sub(r"(?<=[\u4e00-\u9fff])[ \t]+(?=[\u4e00-\u9fff])", "", text)
    residual = NUMBERED_ENTRY_PATTERN.sub("", text)
    return re.sub(r"[\s/·•,;:|]+", "", residual)


def _identity_list_residual(text: str) -> str:
    """Return non-entry content from an already isolated Identity list."""
    text = re.sub(r"(?<=[\u4e00-\u9fff])[ \t]+(?=[\u4e00-\u9fff])", "", text)
    residual = NUMBERED_ENTRY_PATTERN.sub("", text)
    for identity in EXPECTED_META_AGENT_IDENTITIES:
        residual = re.sub(re.escape(identity), "", residual, flags=re.I)
    residual = re.sub(r"(?:元|Agent)", "", residual, flags=re.I)
    return re.sub(r"[\s/·•,;:|]+", "", residual)


def _contract_list_labels(page_number: int) -> tuple[str, ...] | None:
    if page_number == 15:
        return EXPECTED_SKILL_ENTRIES
    if page_number in (7, 14):
        return ("交付官", "业务架构师", "方案架构师", "验证工程师", "审计官")
    return None


def _text_contract_list_text(text: str, page_number: int) -> str:
    """Isolate numbered-entry lines without relying on surrounding copy."""
    labels = _contract_list_labels(page_number)
    if labels is None:
        return ""
    entry_lines: list[str] = []
    for number, label in enumerate(labels, start=1):
        expected = (number, label)
        matches = [
            line
            for line in text.splitlines()
            if expected in _structured_integer_entries(line)
        ]
        if len(matches) != 1:
            return ""
        entry_lines.append(matches[0])
    return "\n".join(entry_lines)


def _pptx_contract_list_text(slide: object, page_number: int) -> str | None:
    """Isolate numbered entries plus nearby standalone cards from their geometry."""
    labels = _contract_list_labels(page_number)
    if labels is None:
        return None

    entry_texts: list[str] = []
    entry_shapes: list[object] = []
    for number, label in enumerate(labels, start=1):
        expected = (number, label)
        matches = [
            shape
            for shape in _iter_shapes(slide.shapes)
            if getattr(shape, "has_text_frame", False)
            and expected in _structured_integer_entries(shape.text)
        ]
        if len(matches) != 1:
            return ""
        entry_shapes.append(matches[0])
        entry_texts.append(matches[0].text)

    left = min(shape.left for shape in entry_shapes)
    right = max(shape.left + shape.width for shape in entry_shapes)
    top = min(shape.top for shape in entry_shapes)
    bottom = max(shape.top + shape.height for shape in entry_shapes)
    horizontal_padding = max(shape.width for shape in entry_shapes)
    vertical_padding = max(shape.height for shape in entry_shapes)
    entry_ids = {shape.shape_id for shape in entry_shapes}
    for shape in _iter_shapes(slide.shapes):
        if (
            not getattr(shape, "has_text_frame", False)
            or not shape.text.strip()
            or shape.shape_id in entry_ids
        ):
            continue
        center_x = shape.left + shape.width / 2
        center_y = shape.top + shape.height / 2
        if (
            left - horizontal_padding <= center_x <= right + horizontal_padding
            and top - vertical_padding <= center_y <= bottom + vertical_padding
        ):
            entry_texts.append(shape.text)
    return "\n".join(entry_texts)


def _pdf_positioned_list_text(page: object, page_number: int) -> str:
    """Extract a list band derived from frozen numbered-entry anchors."""
    spans: list[tuple[float, str, float, float]] = []

    def visitor(text: str, _cm: object, tm: object, _font: object, size: float) -> None:
        if text.strip():
            spans.append((float(tm[4]), float(tm[5]), text, float(size)))

    try:
        page.extract_text(visitor_text=visitor)
    except Exception:  # pragma: no cover - normal extraction is validated elsewhere
        return ""

    if page_number == 15:
        anchor_labels = EXPECTED_SKILL_ENTRIES
        number_pattern = lambda number: rf"^\s*{number}(?=\s|/|$)"
    elif page_number in (7, 14):
        anchor_labels = ("交付官", "业务架构师", "方案架构师", "验证工程师", "审计官")
        number_pattern = lambda number: rf"^\s*0?{number}(?=\s|/|$)"
    else:
        return ""

    anchor_spans: list[tuple[float, float, str, float]] = []
    for number, label in enumerate(anchor_labels, start=1):
        candidates = [
            span
            for span in spans
            if re.match(number_pattern(number), span[2])
            and any(
                abs(span[1] - other[1]) <= 1
                and abs(span[0] - other[0]) <= 120
                and label[:1] in other[2]
                for other in spans
            )
        ]
        if len(candidates) != 1:
            return ""
        anchor_spans.append(candidates[0])
    if len(anchor_spans) != len(anchor_labels):
        return ""

    anchor_ys = [y for _, y, _, _ in anchor_spans]
    padding = max(size for _, _, _, size in anchor_spans) * 2
    low, high = min(anchor_ys) - padding, max(anchor_ys) + padding
    groups: dict[float, list[str]] = {}
    for _, y, text, _ in spans:
        if low <= y <= high:
            groups.setdefault(round(y, 1), []).append(text)
    lines: list[str] = []
    for y in sorted(groups, reverse=True):
        line = ""
        for text in groups[y]:
            if line and re.match(r"\s*\d", text):
                line += " "
            line += text
        lines.append(line)
    return "\n".join(lines)


def _numbered_contract_errors(
    page_label: str,
    page_number: int,
    text: str,
    positioned_list_text: str | None = None,
) -> list[str]:
    errors: list[str] = []
    marker = EXPECTED_PAGE_MARKERS[page_number - 1]
    marker_count = _normalized(text).count(_normalized(marker))
    if marker_count != 1:
        errors.append(
            f"{page_label} must contain exactly one page marker {marker}; found {marker_count}"
        )

    if page_number in (7, 14):
        identity_entries = _structured_integer_entries(text)
        expected_identity_entries = list(
            enumerate(("交付官", "业务架构师", "方案架构师", "验证工程师", "审计官"), start=1)
        )
        if identity_entries != expected_identity_entries:
            errors.append(
                f"{page_label} must contain exactly five fixed Meta Agent entries 01-05; "
                f"found {identity_entries}"
            )
        for identity in EXPECTED_META_AGENT_IDENTITIES:
            if identity not in text:
                errors.append(
                    f"{page_label} is missing fixed Meta Agent identity: "
                    f"{identity}"
                )

    if page_number == 15:
        skill_entries = _structured_integer_entries(text)
        expected_skill_entries = list(enumerate(EXPECTED_SKILL_ENTRIES, start=1))
        if skill_entries != expected_skill_entries:
            errors.append(
                f"{page_label} must contain exactly seven fixed Skill entries 1-7; "
                f"found {skill_entries}"
            )
        for skill in EXPECTED_SKILL_ENTRIES:
            if skill not in text:
                errors.append(
                    f"{page_label} is missing fixed Skill entry: {skill}"
                )
    if page_number == 15:
        list_text = (
            _text_contract_list_text(text, page_number)
            if positioned_list_text is None
            else positioned_list_text
        )
        if list_text and _skill_list_residual(list_text):
            errors.append(
                f"{page_label} contains an unrecognized extra identity/Skill entry"
            )
    elif page_number in (7, 14):
        identity_text = (
            _text_contract_list_text(text, page_number)
            if positioned_list_text is None
            else positioned_list_text
        )
        if identity_text and _identity_list_residual(identity_text):
            errors.append(
                f"{page_label} contains an unrecognized extra identity/Skill entry"
            )
    return errors


def validate(pptx_path: Path, pdf_path: Path | None = None) -> list[str]:
    """Return human-readable validation errors; an empty list means success."""
    errors: list[str] = []
    if not pptx_path.is_file():
        return [f"PPTX file does not exist: {pptx_path}"]

    try:
        presentation = Presentation(pptx_path)
    except Exception as exc:  # pragma: no cover
        return [f"Unable to read PPTX {pptx_path}: {exc}"]

    slide_count = len(presentation.slides)
    if slide_count != EXPECTED_SLIDES:
        errors.append(
            f"PPTX must contain exactly {EXPECTED_SLIDES} slides; found {slide_count}"
        )

    slide_texts = [_slide_text(slide) for slide in presentation.slides]
    for index, (slide, text) in enumerate(
        zip(presentation.slides, slide_texts), start=1
    ):
        if not text.strip():
            errors.append(f"PPTX slide {index} is empty")
        if not _has_visible_native_content(
            slide, presentation.slide_width, presentation.slide_height
        ):
            errors.append(
                f"PPTX slide {index} has no visible native text/content shape"
            )
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                errors.append(f"PPTX slide {index} contains a picture/raster shape")
            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                errors.append(f"PPTX slide {index} contains a media shape")

        normalized_text = _normalized(text)
        errors.extend(
            _numbered_contract_errors(
                "PPTX slide",
                index,
                text,
                _pptx_contract_list_text(slide, index),
            )
        )
        if index <= len(EXPECTED_PAGE_TITLES):
            title = EXPECTED_PAGE_TITLES[index - 1]
            if _normalized(title) not in normalized_text:
                errors.append(
                    f"PPTX slide {index} is missing expected title: {title}"
                )
        for term in SLIDE_REQUIRED_TERMS.get(index, ()):
            if _normalized(term) not in normalized_text:
                errors.append(
                    f"PPTX slide {index} is missing required term: {term}"
                )
        if index <= 4:
            for term in FIRST_PAGE_ML_BAN:
                if _normalized(term) in normalized_text:
                    errors.append(
                        f"PPTX slide {index} uses an ML/NAS term forbidden in the first four pages: {term}"
                    )

    errors.extend(_package_editability_errors(pptx_path))

    deck_text = "\n".join(slide_texts)
    for term in REQUIRED_TERMS:
        if not _required_term_present(deck_text, term):
            errors.append(f"PPTX is missing required term: {term}")
    for declaration in forbidden_declarations(deck_text):
        if declaration in FORBIDDEN_TERMS:
            errors.append(f"PPTX contains forbidden term: {declaration}")
        else:
            errors.append(
                f"PPTX contains forbidden declaration pattern: {declaration}"
            )

    if pdf_path is not None:
        if not pdf_path.is_file():
            errors.append(f"PDF file does not exist: {pdf_path}")
        else:
            try:
                pdf_reader = PdfReader(pdf_path)
                pdf_pages = len(pdf_reader.pages)
            except Exception as exc:  # pragma: no cover
                errors.append(f"Unable to read PDF {pdf_path}: {exc}")
            else:
                if pdf_pages != EXPECTED_SLIDES:
                    errors.append(
                        f"PDF must contain exactly {EXPECTED_SLIDES} pages; found {pdf_pages}"
                    )
                pdf_texts: list[str] = []
                for index, page in enumerate(pdf_reader.pages, start=1):
                    try:
                        page_text = page.extract_text() or ""
                    except Exception as exc:  # pragma: no cover
                        errors.append(f"Unable to extract PDF page {index} text: {exc}")
                        continue
                    pdf_texts.append(page_text)
                    normalized_page = _normalized(page_text)
                    errors.extend(
                        _numbered_contract_errors(
                            "PDF page",
                            index,
                            page_text,
                            _pdf_positioned_list_text(page, index),
                        )
                    )
                    if index <= len(EXPECTED_PAGE_TITLES):
                        title = EXPECTED_PAGE_TITLES[index - 1]
                        if _normalized(title) not in normalized_page:
                            errors.append(
                                f"PDF page {index} is missing expected title: {title}"
                            )
                    for term in SLIDE_REQUIRED_TERMS.get(index, ()):
                        if _normalized(term) not in normalized_page:
                            errors.append(
                                f"PDF page {index} is missing required term: {term}"
                            )
                    if index <= 4:
                        for term in FIRST_PAGE_ML_BAN:
                            if _normalized(term) in normalized_page:
                                errors.append(
                                    f"PDF page {index} uses an ML/NAS term forbidden in the first four pages: {term}"
                                )
                    if index <= len(slide_texts):
                        pptx_text = _normalized(slide_texts[index - 1])
                        if pptx_text and pptx_text not in normalized_page:
                            errors.append(
                                f"PDF page {index} is missing PPTX text: {slide_texts[index - 1][:80]}"
                            )

                pdf_text = "\n".join(pdf_texts)
                for declaration in forbidden_declarations(pdf_text):
                    if declaration in FORBIDDEN_TERMS:
                        errors.append(f"PDF contains forbidden term: {declaration}")
                    else:
                        errors.append(
                            "PDF contains forbidden declaration pattern: "
                            f"{declaration}"
                        )

    return errors


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Validate the final AgentFit submission PPTX and optional PDF."
    )
    parser.add_argument("pptx", type=Path, help="Path to the PPTX file")
    parser.add_argument("pdf", type=Path, nargs="?", help="Optional PDF export")
    args = parser.parse_args()

    errors = validate(args.pptx, args.pdf)
    if errors:
        print("VALIDATION_FAILED")
        for error in errors:
            print(f"- {error}")
        return 1

    presentation = Presentation(args.pptx)
    print(f"pptx_pages={len(presentation.slides)}")
    if args.pdf is not None:
        print(f"pdf_pages={len(PdfReader(args.pdf).pages)}")
    print("content_checks=PASS")
    print("native_editability_checks=PASS")
    if args.pdf is not None:
        print("pdf_page_text_checks=PASS")
    return 0


if __name__ == "__main__":
    sys.exit(main())
