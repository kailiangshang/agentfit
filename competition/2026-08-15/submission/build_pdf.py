#!/usr/bin/env python3
"""Build the submission PDF and contact sheet portably.

The PDF is assembled as two layers per page:
1. an image layer containing the true render of each HTML slide
   (screenshots via a Chromium-family browser in headless mode), and
2. an invisible text layer carrying the PPTX text in shape order, so the
   PDF stays text-extractable and consistent with the editable deck.

Dependencies: python-pptx, pymupdf, pillow. The screenshot browser is
auto-detected (Microsoft Edge / Chrome / Chromium) and overridable via
--browser.
"""

from __future__ import annotations

import argparse
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

import fitz  # pymupdf
from PIL import Image
from pptx import Presentation

ROOT = Path(__file__).resolve().parent
PPTX = ROOT / "agentfit-submission.pptx"
PDF = ROOT / "agentfit-submission.pdf"
SHEET = ROOT / "contact-sheet.jpg"
SLIDES_DIR = ROOT / "slides"

BROWSER_CANDIDATES = (
    "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/usr/bin/chromium",
    "/usr/bin/chromium-browser",
    "/usr/bin/google-chrome",
)

CJK_FONT_CANDIDATES = (
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
)

EMU_W = 12192000
SCALE = 960.0 / EMU_W
HELV_SET = set(chr(c) for c in range(0x20, 0x7F)) | {"·", "…"}


def find_browser(override: str | None) -> str:
    if override:
        return override
    for candidate in BROWSER_CANDIDATES:
        if Path(candidate).is_file():
            return candidate
    found = shutil.which("chromium") or shutil.which("chromium-browser")
    if found:
        return found
    raise SystemExit(
        "no Chromium-family browser found; pass --browser /path/to/browser"
    )


def find_cjk_font() -> str:
    for candidate in CJK_FONT_CANDIDATES:
        if Path(candidate).is_file():
            return candidate
    raise SystemExit("no CJK font found for the text layer")


def screenshot_slides(browser: str, out_dir: Path) -> list[Path]:
    shots = sorted(out_dir.glob("*.png"))
    if shots and len(shots) == 17:
        return shots
    for html in sorted(SLIDES_DIR.glob("[0-9][0-9]-*.html")):
        target = out_dir / f"{html.stem}.png"
        subprocess.run(
            [
                browser,
                "--headless=new",
                "--disable-gpu",
                "--hide-scrollbars",
                "--window-size=1280,720",
                f"--screenshot={target}",
                html.resolve().as_uri(),
            ],
            check=True,
            capture_output=True,
        )
    shots = sorted(out_dir.glob("*.png"))
    if len(shots) != 17:
        raise SystemExit(f"expected 17 screenshots, found {len(shots)}")
    return shots


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from iter_shapes(shapes.shapes)
        else:
            yield shape


def font_runs(line: str):
    buffer, current = [], None
    for ch in line:
        family = "h" if ch in HELV_SET else "c"
        if family != current:
            buffer.append([family, ""])
            current = family
        buffer[-1][1] += ch
    return buffer


def build(pptx_path: Path, pdf_path: Path, sheet_path: Path, browser: str) -> None:
    cjk_font = fitz.Font(fontfile=find_cjk_font())
    helv_font = fitz.Font("helv")

    presentation = Presentation(str(pptx_path))
    with tempfile.TemporaryDirectory(prefix="agentfit-pdf-") as temp:
        shots = screenshot_slides(browser, Path(temp))
        doc = fitz.open()
        for index, slide in enumerate(presentation.slides):
            page = doc.new_page(width=960, height=540)
            page.insert_image(fitz.Rect(0, 0, 960, 540), filename=str(shots[index]))
            writer = fitz.TextWriter(page.rect)
            for shape in iter_shapes(slide.shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text_frame.text
                if not text.strip():
                    continue
                x0 = max(1, shape.left * SCALE) + 1
                base_y = max(10, min(526, shape.top * SCALE + 10))
                for line_index, line in enumerate(text.split("\n")):
                    if not line.strip():
                        continue
                    y = min(534, base_y + line_index * 9)
                    x = x0
                    for family, segment in font_runs(line):
                        font = helv_font if family == "h" else cjk_font
                        writer.append((x, y), segment, font=font, fontsize=7.5)
                        x += font.text_length(segment, fontsize=7.5)
            writer.write_text(page, render_mode=3)
        doc.save(str(pdf_path), garbage=3, deflate=True)

        cols, rows, thumb_w, thumb_h, pad = 4, 5, 320, 180, 8
        sheet = Image.new(
            "RGB",
            (cols * (thumb_w + pad) + pad, rows * (thumb_h + pad) + pad),
            "#f0f0f0",
        )
        for index, page in enumerate(doc):
            pix = page.get_pixmap(
                matrix=fitz.Matrix(thumb_w / 960, thumb_h / 540)
            )
            image = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            if image.size != (thumb_w, thumb_h):
                image = image.resize((thumb_w, thumb_h))
            row, col = divmod(index, cols)
            sheet.paste(
                image, (pad + col * (thumb_w + pad), pad + row * (thumb_h + pad))
            )
        sheet.save(str(sheet_path), quality=88)
    print(f"pdf pages: 17 -> {pdf_path}")
    print(f"contact sheet -> {sheet_path}")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--pptx", type=Path, default=PPTX)
    parser.add_argument("--pdf", type=Path, default=PDF)
    parser.add_argument("--sheet", type=Path, default=SHEET)
    parser.add_argument("--browser", default=None)
    args = parser.parse_args()
    build(args.pptx, args.pdf, args.sheet, find_browser(args.browser))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
