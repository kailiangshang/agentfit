#!/usr/bin/env python3
"""Contract tests for the final AgentFit preliminary submission."""

from __future__ import annotations

import base64
import html
import importlib.util
import re
import subprocess
import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen.canvas import Canvas


ROOT = Path(__file__).resolve().parent
COMPETITION_ROOT = ROOT.parent
REPO_ROOT = ROOT.parents[2]
SLIDES_DIR = ROOT / "slides"
INTRODUCTION = ROOT / "work-introduction.md"
OUTLINE = ROOT / "ppt-outline.md"
VALIDATOR = ROOT / "validate_presentation.py"
BUILDER = ROOT / "build_presentation.py"
PPTX = ROOT / "agentfit-submission.pptx"
PDF = ROOT / "agentfit-submission.pdf"
AGENT_IDENTITY = ROOT / "agent-identity.md"
SKILL_CATALOG = ROOT / "skill-catalog.md"
RISK_GATES = ROOT / "risk-and-human-gates.md"
OPENNESS = ROOT / "openness-and-compliance.md"
SOLUTION = REPO_ROOT / "docs/agentfit-solution.md"
HOME_DEMO_RUNBOOK = REPO_ROOT / "docs/guides/home-demo-runbook.md"
EVIDENCE_REGISTRY = REPO_ROOT / "docs/internal/evidence-research/evidence-registry.json"
OPSPILOT_EVIDENCE_CARD = (
    REPO_ROOT
    / "docs/internal/evidence-research/cards/operations/opspilot-zero-demo.md"
)

EVALUATION_IDENTITY = "CandidateVersion × SampleVersion × RunIndex"


def _load_validator():
    spec = importlib.util.spec_from_file_location("validate_presentation", VALIDATOR)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"unable to import {VALIDATOR}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _load_builder():
    spec = importlib.util.spec_from_file_location("build_presentation", BUILDER)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"unable to import {BUILDER}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _introduction_body(text: str) -> str:
    match = re.search(
        r"^## 500 字以内作品简介\s*$\n(?P<body>.*?)(?=^##\s|\Z)",
        text,
        flags=re.MULTILINE | re.DOTALL,
    )
    if not match:
        raise AssertionError("missing '500 字以内作品简介' section")
    return match.group("body").strip()


def _pdf_layout_lines(page: int) -> list[str]:
    completed = subprocess.run(
        [
            "pdftotext",
            "-f",
            str(page),
            "-l",
            str(page),
            "-layout",
            str(PDF),
            "-",
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    return [re.sub(r"\s+", " ", line).strip() for line in completed.stdout.splitlines()]


def _forward_arrows_between(slide: object, source: object, target: object) -> list[object]:
    return [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.strip() == "→"
        and source.left + source.width <= shape.left
        and shape.left + shape.width <= target.left
    ]


def _html_title(text: str) -> str:
    match = re.search(
        r'<h1\s+class="title"[^>]*>(?P<title>.*?)</h1>',
        text,
        flags=re.DOTALL,
    )
    if match is None:
        raise AssertionError("missing h1.title")
    title = re.sub(r"<[^>]+>", "", match.group("title"))
    return re.sub(r"\s+", "", html.unescape(title))


class SubmissionContractTest(unittest.TestCase):
    def test_repository_has_one_final_submission_and_no_parallel_versions(self) -> None:
        self.assertEqual("submission", ROOT.name)
        self.assertFalse((COMPETITION_ROOT / "alternatives").exists())
        self.assertEqual(
            [ROOT / "agentfit-submission.pptx"],
            sorted(COMPETITION_ROOT.rglob("*.pptx")),
        )
        self.assertEqual(
            [ROOT / "agentfit-submission.pdf"],
            sorted(COMPETITION_ROOT.rglob("*.pdf")),
        )
        self.assertEqual(
            [REPO_ROOT / "docs/agentfit-solution.md"],
            sorted(REPO_ROOT.glob("docs/*solution*.md")),
        )

    def test_final_submission_keeps_required_detailed_contracts(self) -> None:
        for filename in (
            "agent-identity.md",
            "skill-catalog.md",
            "risk-and-human-gates.md",
            "openness-and-compliance.md",
        ):
            with self.subTest(filename=filename):
                self.assertTrue((ROOT / filename).is_file())

    def test_audience_materials_avoid_rubric_answering_language(self) -> None:
        audience_materials = (
            AGENT_IDENTITY,
            SKILL_CATALOG,
            RISK_GATES,
            OPENNESS,
            OUTLINE,
            ROOT / "README.md",
            INTRODUCTION,
            *sorted(SLIDES_DIR.glob("*.html")),
        )
        forbidden = (
            "官方字段对照",
            "官方要求",
            "与官方要求的映射",
            "本清单按此",
            "本清单如何满足",
            "覆盖要求",
            "不按数量凑分",
            "红线对齐",
            "官方 8 字段",
            "Skill 必选",
            "上下文 4 选 2",
            "上下文能力 · 4 选 2",
            "8 字段契约",
            "10 字段契约",
        )
        for path in audience_materials:
            text = path.read_text(encoding="utf-8")
            for term in forbidden:
                with self.subTest(path=path, term=term):
                    self.assertNotIn(term, text)

    def test_detailed_agent_identity_contract_is_complete(self) -> None:
        text = AGENT_IDENTITY.read_text(encoding="utf-8")
        for agent in (
            "EngagementLead",
            "BusinessEngineer",
            "AgentArchitect",
            "ValidationEngineer",
            "GovernanceAuditor",
        ):
            with self.subTest(agent=agent):
                self.assertIn(agent, text)
        for field in (
            "Name",
            "Role",
            "Capabilities",
            "Inputs",
            "Outputs",
            "Dependencies",
            "Decision Boundary",
            "Trace",
        ):
            with self.subTest(field=field):
                self.assertIn(f"| {field} |", text)
        self.assertIn("真实实例 `NOT_STARTED`", text)
        self.assertIn("sealed_holdout_outcome_consumer: GovernanceAuditor only", text)

    def test_detailed_skill_contract_is_complete(self) -> None:
        text = SKILL_CATALOG.read_text(encoding="utf-8")
        for skill_id in range(1, 8):
            with self.subTest(skill_id=skill_id):
                self.assertIn(f"### S{skill_id}.", text)
        for field in (
            "名称",
            "类型",
            "使用场景",
            "输入参数",
            "输出结果",
            "调用条件",
            "依赖工具系统",
            "失败处理",
            "权限安全",
            "复用价值",
        ):
            with self.subTest(field=field):
                self.assertGreaterEqual(text.count(f"| {field} |"), 7)
        self.assertIn("真实绑定 `NOT_STARTED`", text)

    def test_detailed_risk_and_human_gate_contract_is_complete(self) -> None:
        text = RISK_GATES.read_text(encoding="utf-8")
        for term in (
            "Human 门禁",
            "审批主体",
            "拒绝路径",
            "回滚路径",
            "Trace 要求",
            "高风险动作无批准",
            "真实运行状态为 `NOT_STARTED`",
        ):
            with self.subTest(term=term):
                self.assertIn(term, text)

    def test_openness_contract_discloses_current_and_conditional_state(self) -> None:
        text = OPENNESS.read_text(encoding="utf-8")
        for term in (
            "仓库当前为私有",
            "尚未发布",
            "后续工作不会自动启动",
            "晋级结果",
            "明确授权",
            "附录 A5",
        ):
            with self.subTest(term=term):
                self.assertIn(term, text)
        self.assertNotIn("路演第 9 页", text)

    def test_opspilot_baseline_has_minimal_source_evidence_card(self) -> None:
        self.assertTrue(OPSPILOT_EVIDENCE_CARD.is_file())
        card = OPSPILOT_EVIDENCE_CARD.read_text(encoding="utf-8")
        for term in (
            "https://assets.datawhale.cn/131266/dashboard/1785575974456/opspilot-zero-demo.zip",
            "0bb0f37c227fb5031cd66b6d69dbcbd533602c26b7d5e93f66f93fa02f653478",
            "4 个 Agent",
            "7 个 Skill",
            "db_pool_exhausted",
            "slow_sql_degradation",
            "2026-08-12",
            "不证明 AgentFit 已运行",
        ):
            with self.subTest(term=term):
                self.assertIn(term, card)

        import json

        registry = json.loads(EVIDENCE_REGISTRY.read_text(encoding="utf-8"))
        entries = {
            entry["evidence_id"]: entry for entry in registry.get("entries", [])
        }
        self.assertIn("opspilot-zero-demo", entries)
        self.assertEqual(
            "docs/internal/evidence-research/cards/operations/opspilot-zero-demo.md",
            entries["opspilot-zero-demo"]["card_path"],
        )

    def test_documents_separate_preliminary_and_later_stages(self) -> None:
        documents = (
            COMPETITION_ROOT / "README.md",
            REPO_ROOT / "docs/README.md",
            REPO_ROOT / "docs/agentfit-solution.md",
        )
        for path in documents:
            text = path.read_text(encoding="utf-8")
            with self.subTest(path=path):
                self.assertIn("2026-08-15 初赛提交阶段", text)
                self.assertIn("后续阶段", text)
                self.assertIn("晋级结果", text)

    def test_home_demo_runbook_is_executable_and_evidence_safe(self) -> None:
        self.assertTrue(HOME_DEMO_RUNBOOK.is_file())
        runbook = HOME_DEMO_RUNBOOK.read_text(encoding="utf-8")
        docs_index = (REPO_ROOT / "docs/README.md").read_text(encoding="utf-8")

        for term in (
            "https://github.com/sierra-research/tau2-bench",
            "v1.0.1",
            "retail",
            "114",
            "--task-ids 0",
            "--task-ids 0 1 2",
            "data/simulations/",
            "EngagementLead",
            "BusinessEngineer",
            "AgentArchitect",
            "ValidationEngineer",
            "GovernanceAuditor",
            "ProjectCase != Sample",
            "AgentFit adaptation/synthetic",
            "M0",
            "M1",
            "IN_PROGRESS",
            "GovernanceAuditor only",
            "Agentless、单 Agent 和多 Agent",
            ".env",
            "不得提交",
            'importlib.metadata as m; print(m.version("tau2"))',
            "preflight-only",
            "不是 AgentFit Candidate",
            "m0-authorization.md",
            "agentteams-version.txt",
            "agentteams-status.json",
            "m0-baseline.md",
            "known_boundaries",
            "source/task-0.json",
            "2>&1 | tee",
            "agentfit-retail-preflight-12/results.json",
            "agentfit-retail-preflight-20/results.json",
        ):
            with self.subTest(term=term):
                self.assertIn(term, runbook)

        self.assertNotIn("tau2 --version", runbook)
        self.assertNotIn("已执行候选证据", runbook)
        self.assertIn('AGENTFIT_ROOT="$(git rev-parse --show-toplevel)"', runbook)
        self.assertIn('AGENTFIT_TAU3_ROOT="$AGENTFIT_ROOT/../agentfit-labs/tau2-bench"', runbook)
        self.assertNotIn("cd ../agentfit-labs/tau2-bench", runbook)
        self.assertIn("export AGENTFIT_AGENT_MODEL", runbook)
        self.assertIn('test -n "$AGENTFIT_AGENT_MODEL"', runbook)
        self.assertIn(
            'jq \'map(select(.id == "0")) | .[0]\' "${AGENTFIT_TAU3_ROOT}/data/tau2/domains/retail/tasks.json" > "${AGENTFIT_RUN_ROOT}/source/task-0.json"',
            runbook,
        )
        self.assertIn(
            'git -C "$AGENTFIT_TAU3_ROOT" check-ignore -q .env',
            runbook,
        )
        self.assertIn("公开 `test` split 不是 sealed holdout", runbook)
        self.assertIn("OpsPilot 与 retail 保持两个独立来源", runbook)
        self.assertIn("M2/M3/M4 均未启动", runbook)
        for size in (12, 20):
            result_dir = f"agentfit-retail-preflight-{size}"
            self.assertIn(
                f'test -f "$AGENTFIT_TAU3_ROOT/data/simulations/{result_dir}/results.json"',
                runbook,
            )
            self.assertIn(
                f'cp -a "$AGENTFIT_TAU3_ROOT/data/simulations/{result_dir}" "$AGENTFIT_RUN_ROOT/native-runs/"',
                runbook,
            )

        self.assertIn("[回家 Demo 执行手册](guides/home-demo-runbook.md)", docs_index)

    def test_home_demo_runtime_outputs_are_local_only(self) -> None:
        gitignore = (REPO_ROOT / ".gitignore").read_text(encoding="utf-8")
        self.assertIn(".local-demo/", gitignore.splitlines())

    def test_canonical_solution_guides_later_implementation_without_contract_drift(self) -> None:
        solution = SOLUTION.read_text(encoding="utf-8")
        identity = AGENT_IDENTITY.read_text(encoding="utf-8")
        skills = SKILL_CATALOG.read_text(encoding="utf-8")

        self.assertEqual(1, solution.count("SampleSemanticSpec = {"))
        self.assertEqual(
            1,
            solution.count("核心能力缺口无法在授权范围内补齐"),
        )
        self.assertIn(
            "单项目必须包含 adaptation、validation、sealed holdout 和 "
            "stress and failure 四份互异的 SampleSetManifest",
            solution,
        )
        self.assertIn("### 13.3 后续最小实施顺序与阶段完成定义", solution)
        for milestone in (
            "M0 · 启动授权与基线冻结",
            "M1 · 手动可审计 walking skeleton",
            "M2 · 确定性合同代码化",
            "M3 · 统一候选对照",
            "M4 · 复现与比赛证据包",
        ):
            with self.subTest(milestone=milestone):
                self.assertIn(milestone, solution)
        self.assertIn("§4.10(Agent 严格定义)", identity)
        self.assertIn("§4.9(能力语义)", skills)
        self.assertIn("§4.11(任务—能力对齐)", skills)
        self.assertNotIn("`AUTHORIZED`", solution)
        self.assertIn(
            "SampleSemanticSpec、Sample、SampleSetManifest",
            OPENNESS.read_text(encoding="utf-8"),
        )
        self.assertIn(
            "以下顺序只在第 13.2 节第一段的外部授权条件满足后生效",
            solution,
        )
        self.assertIn(
            "M1 只证明当前环境首次贯通，不构成可复现最小闭环完成声明",
            solution,
        )
        self.assertIn(
            "独立复现成功且第 13.2 节七项完成门禁全部满足后，才可声明"
            "“AgentFit 已在 AgentTeams 跑通最小闭环”",
            solution,
        )
        self.assertIn(
            "必须真实运行 Agentless、单 Agent 和多 Agent 三类候选",
            solution,
        )
        self.assertIn(
            "Human 混合候选必须真实运行，或由 GovernanceAuditor 记录不适用理由、"
            "证据与重新评估条件",
            solution,
        )
        self.assertEqual(
            2,
            solution.count(
                "候选冻结后只有 GovernanceAuditor 可以解析 sealed holdout"
            ),
        )

    def test_deck_keeps_later_runtime_work_conditional(self) -> None:
        conclusion = (SLIDES_DIR / "12-conclusion.html").read_text(encoding="utf-8")
        self.assertIn("若获准，后续门禁", conclusion)
        self.assertNotIn("下一步：冻结 ProjectCase", conclusion)

    def test_agentteams_and_agentfit_ownership_are_explicit(self) -> None:
        expected = "AgentTeams 承载 Worker、Team、Room、Human；AgentFit 落地 Dossier 与 Trace。"
        for path in (
            SLIDES_DIR / "08-agentteams.html",
            ROOT / "README.md",
            OUTLINE,
            VALIDATOR,
        ):
            with self.subTest(path=path):
                self.assertIn(expected, path.read_text(encoding="utf-8"))

    def test_historical_process_documents_are_removed(self) -> None:
        for path in (
            COMPETITION_ROOT / "design",
            COMPETITION_ROOT / "planning",
            COMPETITION_ROOT / "research",
            REPO_ROOT / "docs/superpowers",
            REPO_ROOT / ("_render_" + "fusion.py"),
            REPO_ROOT / ".superpowers",
        ):
            with self.subTest(path=path):
                self.assertFalse(path.exists())

    def test_final_submission_directory_is_self_contained(self) -> None:
        self.assertTrue(SLIDES_DIR.is_dir())
        self.assertTrue(VALIDATOR.is_file())
        self.assertTrue(BUILDER.is_file())
        self.assertTrue(INTRODUCTION.is_file())
        self.assertTrue(OUTLINE.is_file())
        root = str(ROOT)
        for path in ROOT.rglob("*"):
            self.assertTrue(
                str(path).startswith(root),
                f"stray file outside final submission: {path}",
            )

    def test_builder_runs_from_repository_root(self) -> None:
        builder = _load_builder()
        self.assertEqual(REPO_ROOT, builder.REPO_ROOT)

    def test_final_sources_contain_no_superseded_version_narrative(self) -> None:
        forbidden = (
            "fusion-v3",
            "scoreline-v2",
            "alternatives/",
            "第三套替代",
            "替代路演候选",
            "融合版",
            "评分主线版",
            "冻结版",
            "agentfit-preliminary-draft",
            "agentfit-fusion-v3",
            "test_fusion_contract",
            "FusionV3ContractTest",
        )
        paths = (
            list(ROOT.glob("*.md"))
            + list(ROOT.glob("*.py"))
            + list(SLIDES_DIR.glob("*.html"))
            + [COMPETITION_ROOT / "README.md", REPO_ROOT / "docs/README.md"]
        )
        for path in paths:
            if path == Path(__file__):
                continue
            text = path.read_text(encoding="utf-8")
            for term in forbidden:
                with self.subTest(path=path, term=term):
                    self.assertNotIn(term, text)

    def test_canonical_markdown_relative_links_resolve(self) -> None:
        markdown_files = (
            list(ROOT.glob("*.md"))
            + [
                COMPETITION_ROOT / "README.md",
                REPO_ROOT / "docs/README.md",
                REPO_ROOT / "docs/agentfit-solution.md",
                HOME_DEMO_RUNBOOK,
            ]
        )
        link_pattern = re.compile(r"\[[^\]]+\]\((?!https?://|#)([^)#]+)(?:#[^)]+)?\)")
        for path in markdown_files:
            text = path.read_text(encoding="utf-8")
            for raw_target in link_pattern.findall(text):
                target = (path.parent / raw_target).resolve()
                with self.subTest(path=path, target=raw_target):
                    self.assertTrue(target.exists(), f"broken link: {path} -> {raw_target}")

    def test_seventeen_html_slides_with_stable_names(self) -> None:
        slides = sorted(SLIDES_DIR.glob("[0-9][0-9]-*.html"))
        self.assertEqual(17, len(slides))
        self.assertEqual(
            [f"{index:02d}" for index in range(1, 18)],
            [path.name[:2] for path in slides],
        )

    def test_submission_page_and_contract_counts_are_exact(self) -> None:
        self.assertEqual(17, len(Presentation(PPTX).slides))
        self.assertEqual(17, len(PdfReader(PDF).pages))

        identity = AGENT_IDENTITY.read_text(encoding="utf-8")
        identity_headings = re.findall(
            r"^## (\d+)\.\s", identity, flags=re.MULTILINE
        )
        self.assertEqual(["1", "2", "3", "4", "5"], identity_headings)

        skills = SKILL_CATALOG.read_text(encoding="utf-8")
        skill_headings = re.findall(
            r"^### (S\d+)\.\s", skills, flags=re.MULTILINE
        )
        self.assertEqual([f"S{index}" for index in range(1, 8)], skill_headings)

    def test_html_sources_have_exact_markers_and_titles(self) -> None:
        validator = _load_validator()
        slides = sorted(SLIDES_DIR.glob("[0-9][0-9]-*.html"))
        for index, path in enumerate(slides, start=1):
            source = path.read_text(encoding="utf-8")
            marker = validator.EXPECTED_PAGE_MARKERS[index - 1]
            with self.subTest(page=index, marker=marker):
                self.assertEqual(1, source.count(marker))
            with self.subTest(page=index, title=validator.EXPECTED_PAGE_TITLES[index - 1]):
                self.assertEqual(
                    re.sub(r"\s+", "", validator.EXPECTED_PAGE_TITLES[index - 1]),
                    _html_title(source),
                )

    def test_delivered_artifacts_have_exact_identity_and_skill_entries(self) -> None:
        validator = _load_validator()
        presentation = Presentation(PPTX)
        pdf_reader = PdfReader(PDF)
        for label, text in (
            ("PPTX slide 7", validator._slide_text(presentation.slides[6])),
            ("PPTX slide 14", validator._slide_text(presentation.slides[13])),
            ("PDF page 7", pdf_reader.pages[6].extract_text() or ""),
            ("PDF page 14", pdf_reader.pages[13].extract_text() or ""),
        ):
            with self.subTest(artifact=label):
                self.assertEqual(
                    [
                        (1, "交付官"),
                        (2, "业务架构师"),
                        (3, "方案架构师"),
                        (4, "验证工程师"),
                        (5, "审计官"),
                    ],
                    validator._structured_integer_entries(text),
                )
                for identity in validator.EXPECTED_META_AGENT_IDENTITIES:
                    self.assertIn(identity, text)

        for label, text in (
            ("PPTX slide 15", validator._slide_text(presentation.slides[14])),
            ("PDF page 15", pdf_reader.pages[14].extract_text() or ""),
        ):
            with self.subTest(artifact=label):
                self.assertEqual(
                    list(enumerate(validator.EXPECTED_SKILL_ENTRIES, start=1)),
                    validator._structured_integer_entries(text),
                )
                for skill in validator.EXPECTED_SKILL_ENTRIES:
                    self.assertIn(skill, text)

    def test_revised_narrative_and_evidence_boundaries_are_frozen(self) -> None:
        page_5 = (SLIDES_DIR / "05-search-space.html").read_text(encoding="utf-8")
        page_6 = (SLIDES_DIR / "06-selection-rule.html").read_text(encoding="utf-8")
        page_11 = (SLIDES_DIR / "11-evidence.html").read_text(encoding="utf-8")

        for term in (
            "完整方案空间",
            "工具",
            "Skill",
            "MCP",
            "Memory",
            "模型",
            "Agent 拓扑",
            "Human 边界",
            "C0",
            "C1",
            "C2",
            "C3",
        ):
            with self.subTest(page=5, term=term):
                self.assertIn(term, page_5)
        for term in (
            "五阶段闭环",
            "定义案例与验收",
            "构建最小候选",
            "运行并测量",
            "分析并调整",
            "验证并停止",
        ):
            with self.subTest(page=6, term=term):
                self.assertIn(term, page_6)
        self.assertIn(
            "按失败模式调整完整方案七维，不越过 Human 门禁",
            page_6,
        )
        for term in (
            "Tool",
            "Skill",
            "MCP",
            "Memory",
            "模型",
            "Agent 拓扑",
            "Human 边界",
        ):
            with self.subTest(page=6, complete_solution_dimension=term):
                self.assertIn(term, page_6)
        validator = _load_validator()
        self.assertIn(validator.SLIDE_11_EVIDENCE_STATEMENT, page_11)

        source = "\n".join(
            path.read_text(encoding="utf-8")
            for path in sorted(SLIDES_DIR.glob("*.html"))
        )
        self.assertEqual([], validator.forbidden_declarations(source))

    def test_html_sources_reject_normalized_forbidden_claim_mutations(self) -> None:
        validator = _load_validator()
        source = (SLIDES_DIR / "17-a5-openness.html").read_text(encoding="utf-8")
        for claim in (
            "AutoML",
            "已验证 Meta-learning",
            "Meta-learning is verified",
            "META-LEARNING IS VERIFIED",
            "Meta learning is verified",
            "Meta-learning has been verified",
            "verified Meta learning",
            "VERIFIED META LEARNING",
            "AgentTeams end-to-end integration complete",
        ):
            with self.subTest(claim=claim):
                declarations = validator.forbidden_declarations(source + claim)
                self.assertTrue(declarations, claim)

    def test_html_sources_use_pptx_safe_solid_colors(self) -> None:
        for path in sorted(SLIDES_DIR.glob("*.html")):
            source = path.read_text(encoding="utf-8")
            with self.subTest(slide=path.name):
                self.assertNotRegex(source, r"\b(?:rgba|hsla)\s*\(")

    def test_final_convergence_copy_and_diagram_contract(self) -> None:
        cover = (SLIDES_DIR / "01-cover.html").read_text(encoding="utf-8")
        selection = (SLIDES_DIR / "06-selection-rule.html").read_text(
            encoding="utf-8"
        )
        search = (SLIDES_DIR / "13-a1-search.html").read_text(encoding="utf-8")
        risk = (SLIDES_DIR / "16-a4-risk.html").read_text(encoding="utf-8")
        agentteams = (SLIDES_DIR / "08-agentteams.html").read_text(encoding="utf-8")
        evidence = (SLIDES_DIR / "11-evidence.html").read_text(encoding="utf-8")

        self.assertIn("C2 · N 业务 Agent", cover)

        self.assertIn("候选搜索顺序 · 设计契约（非运行结果）", selection)
        for label in (
            "C0 · Agentless · 待真实试验",
            "C1 · 单 Agent · 待真实试验",
            "C2 · 多 Agent · 待真实试验",
            "C3 · Human 混合 · 待真实试验",
        ):
            with self.subTest(label=label):
                self.assertIn(label, selection)
        for misleading_mark in ("✓", "⚠", "验收阈值", "非真实数据"):
            with self.subTest(misleading_mark=misleading_mark):
                self.assertNotIn(misleading_mark, selection)

        self.assertIn("生成 → 定位 · 最多 3 次 · 未通过退出", search)
        self.assertIn("固定 G / Π / ρ；只更新 θ 与分区内状态", search)
        self.assertIn("更新 G / Π / ρ；比较 validation samples", search)

        self.assertIn("六类异常进入同一 Trace 语义：", risk)
        self.assertIn("对应恢复：", risk)
        self.assertIn("后续运行映射", agentteams)
        self.assertIn("若获准，下一门禁", evidence)
        self.assertNotIn("第一阶段实现", agentteams)

    def test_final_convergence_pdf_labels_do_not_wrap(self) -> None:
        expected_by_page = {
            1: ("C2 · N 业务 Agent",),
            2: (
                "AGENTFIT 拿它做什么",
                "demo 事故 · db.pool.maxSize: 50 -> 8",
            ),
            6: (
                "候选搜索顺序 · 设计契约（非运行结果）",
                "C0 · Agentless · 待真实试验",
                "C1 · 单 Agent · 待真实试验",
                "C2 · 多 Agent · 待真实试验",
                "C3 · Human 混合 · 待真实试验",
            ),
            7: ("作出交付决定",),
            11: (
                "READY · 已完成",
                "SMOKE · 平台已单独试用",
                "NOT STARTED · 仍待运行",
            ),
            16: (
                "异常与恢复 · 同一 Trace 语义",
                "六类异常进入同一 Trace 语义：",
            ),
        }
        for page, labels in expected_by_page.items():
            lines = _pdf_layout_lines(page)
            for label in labels:
                with self.subTest(page=page, label=label):
                    self.assertTrue(any(label in line for line in lines), lines)

        page_13 = _pdf_layout_lines(13)
        self.assertFalse(
            any(line == "据。" for line in page_13),
            f"page 13 contains an orphaned final character: {page_13}",
        )

    def test_final_candidate_graph_has_no_text_overlap(self) -> None:
        presentation = Presentation(PPTX)
        slide = presentation.slides[12]
        by_text = {
            shape.text.strip(): shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        }
        arrow = by_text["←"]
        exit_condition = by_text["生成 → 定位 · 最多 3 次 · 未通过退出"]
        self.assertLessEqual(arrow.top + arrow.height, exit_condition.top)

    def test_final_candidate_graph_and_blocks_geometry(self) -> None:
        presentation = Presentation(PPTX)

        selection_slide = presentation.slides[5]
        candidate_texts = (
            "C0 · Agentless · 待真实试验",
            "C1 · 单 Agent · 待真实试验",
            "C2 · 多 Agent · 待真实试验",
            "C3 · Human 混合 · 待真实试验",
        )
        candidate_text_shapes = [
            next(shape for shape in selection_slide.shapes if shape.text == text)
            for text in candidate_texts
        ]
        candidate_blocks = [
            next(
                shape
                for shape in selection_slide.shapes
                if not getattr(shape, "text", "").strip()
                and shape.left <= text_shape.left
                and text_shape.left + text_shape.width <= shape.left + shape.width
                and shape.top <= text_shape.top
                and text_shape.top + text_shape.height <= shape.top + shape.height
                and shape.width < Inches(3)
            )
            for text_shape in candidate_text_shapes
        ]
        self.assertEqual(1, len({shape.width for shape in candidate_blocks}))
        self.assertEqual(1, len({shape.height for shape in candidate_blocks}))
        self.assertEqual(1, len({shape.top for shape in candidate_blocks}))

        search_slide = presentation.slides[12]
        by_text = {
            shape.text.strip(): shape
            for shape in search_slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        }
        forward_edges = (
            ("去重", "检索"),
            ("检索", "定位"),
            ("定位", "生成"),
            ("生成", "测试"),
            ("测试", "人工"),
        )
        for source, target in forward_edges:
            self.assertLess(by_text[source].left, by_text[target].left)
            arrows = _forward_arrows_between(
                search_slide, by_text[source], by_text[target]
            )
            self.assertEqual(
                1,
                len(arrows),
                f"expected exactly one forward arrow for {source} → {target}",
            )

            mutated = Presentation(PPTX)
            mutated_slide = mutated.slides[12]
            mutated_by_text = {
                shape.text.strip(): shape
                for shape in mutated_slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            }
            mutated_arrow = _forward_arrows_between(
                mutated_slide, mutated_by_text[source], mutated_by_text[target]
            )[0]
            mutated_arrow._element.getparent().remove(mutated_arrow._element)
            self.assertEqual(
                [],
                _forward_arrows_between(
                    mutated_slide, mutated_by_text[source], mutated_by_text[target]
                ),
                f"mutation proof failed for removed edge {source} → {target}",
            )

        locator = by_text["定位"]
        generator = by_text["生成"]

        return_arrow = by_text["←"]
        return_line = next(
            shape
            for shape in search_slide.shapes
            if not getattr(shape, "text", "").strip()
            and shape.height <= Inches(0.03)
            and locator.left < shape.left
            and shape.left < generator.left + generator.width
            and shape.top > locator.top + locator.height
        )
        self.assertLess(return_arrow.left, return_line.left)
        self.assertGreater(return_line.width, Inches(1.0))

    def test_introduction_under_five_hundred_non_whitespace_chars(self) -> None:
        body = _introduction_body(INTRODUCTION.read_text(encoding="utf-8"))
        count = len(re.sub(r"\s+", "", body))
        self.assertLessEqual(count, 500)
        for term in ("AgentFit", "AgentTeams", "OpsPilot", "方案建筑师"):
            self.assertIn(term, body)

    def test_first_four_slides_avoid_ml_and_nas_terminology(self) -> None:
        banned = (
            "Architecture Search",
            "NAS",
            "Meta-learning",
            "Meta learning",
            "七层",
            "L1",
            "L2",
            "L3",
            "SVD",
            "argmin",
            "Pareto",
            "贝叶斯",
            "inner loop",
            "outer loop",
        )
        for index in range(1, 5):
            matches = sorted(SLIDES_DIR.glob(f"{index:02d}-*.html"))
            self.assertEqual(1, len(matches), f"expected one slide for page {index}")
            text = matches[0].read_text(encoding="utf-8")
            for term in banned:
                with self.subTest(page=index, term=term):
                    self.assertNotIn(term, text)

    def test_cover_keeps_general_product_position_not_ops_product(self) -> None:
        cover = (SLIDES_DIR / "01-cover.html").read_text(encoding="utf-8")
        for term in ("方案建筑师", "AgentTeams"):
            self.assertIn(term, cover)
        for term in ("运维产品", "OpsPilot Zero 是 AgentFit"):
            self.assertNotIn(term, cover)

    def test_meta_agent_vs_business_agent_two_layers_are_explicit(self) -> None:
        meta_team = (SLIDES_DIR / "07-meta-team.html").read_text(encoding="utf-8")
        for term in (
            "EngagementLead",
            "BusinessEngineer",
            "AgentArchitect",
            "ValidationEngineer",
            "GovernanceAuditor",
            "元 Agent",
            "业务执行 Agent",
        ):
            self.assertIn(term, meta_team)
        self.assertIn("共同覆盖闭环", meta_team)
        self.assertIn("交接责任", meta_team)

    def test_compact_pdf_labels_do_not_wrap_orphan_characters(self) -> None:
        page_7 = _pdf_layout_lines(7)
        for label in (
            "01 交付官",
            "02 业务架构师",
            "03 方案架构师",
            "04 验证工程师",
            "05 审计官",
        ):
            with self.subTest(page=7, label=label):
                self.assertTrue(any(label in line for line in page_7), page_7)

        page_9 = _pdf_layout_lines(9)
        for label in (
            "共享状态 · 项目事实源",
            "示例：读配置 · 列变更 · 写回滚",
        ):
            with self.subTest(page=9, label=label):
                self.assertTrue(any(label in line for line in page_9), page_9)

        page_17 = _pdf_layout_lines(17)
        self.assertTrue(
            any("baseline 引用与核验" in line for line in page_17),
            page_17,
        )

    def test_baseline_incidents_and_task_sample_terms_are_present(self) -> None:
        source = "\n".join(
            path.read_text(encoding="utf-8")
            for path in sorted(SLIDES_DIR.glob("*.html"))
        )
        for term in (
            "OpsPilot",
            "ProjectCase",
            "db_pool_exhausted",
            "slow_sql_degradation",
            "TaskSample",
            "Episode",
            EVALUATION_IDENTITY,
            "设计契约，非运行证据",
            "requires_runtime_trial",
            "同一搜索空间",
        ):
            with self.subTest(term=term):
                self.assertIn(term, source)

    def test_search_space_and_seven_layer_mapping_restored(self) -> None:
        source = "\n".join(
            path.read_text(encoding="utf-8")
            for path in sorted(SLIDES_DIR.glob("*.html"))
        )
        for term in (
            "七层 ML 映射",
            "G, Π, θ, ρ",
            "inner loop",
            "outer loop",
            "Meta-learning",
        ):
            with self.subTest(term=term):
                self.assertIn(term, source)

    def test_no_fabricated_runtime_evidence_in_html(self) -> None:
        source = "\n".join(
            path.read_text(encoding="utf-8")
            for path in sorted(SLIDES_DIR.glob("*.html"))
        )
        forbidden = (
            "已选定 C2",
            "已选定 C1",
            "已跑通最小闭环",
            "C2 胜出",
            "ROI 提升",
            "准确率 9",
            "ImageNet",
            "90%+",
            "已开放",
            "Meta-learning 已验证",
        )
        for term in forbidden:
            with self.subTest(term=term):
                self.assertNotIn(term, source)

    def test_validator_encodes_submission_contract(self) -> None:
        validator = _load_validator()
        self.assertEqual(17, validator.EXPECTED_SLIDES)
        self.assertTrue(hasattr(validator, "EXPECTED_PAGE_TITLES"))
        self.assertEqual(17, len(validator.EXPECTED_PAGE_TITLES))
        self.assertEqual(17, len(validator.EXPECTED_PAGE_MARKERS))
        self.assertEqual(
            [f"{index:02d} / 12" for index in range(1, 13)]
            + [f"A{index} / 05" for index in range(1, 6)],
            list(validator.EXPECTED_PAGE_MARKERS),
        )
        self.assertEqual(
            (
                "EngagementLead",
                "BusinessEngineer",
                "AgentArchitect",
                "ValidationEngineer",
                "GovernanceAuditor",
            ),
            validator.EXPECTED_META_AGENT_IDENTITIES,
        )
        self.assertEqual(
            ("任务编译", "能力对齐", "候选建图", "统一试验", "独立审计", "人工门禁", "经验沉淀"),
            validator.EXPECTED_SKILL_ENTRIES,
        )
        self.assertIn("OpsPilot", validator.SLIDE_11_EVIDENCE_STATEMENT)
        self.assertIn("官方案例锚点", validator.SLIDE_11_EVIDENCE_STATEMENT)
        self.assertIn("retail / airline", validator.SLIDE_11_EVIDENCE_STATEMENT)
        self.assertTrue(hasattr(validator, "SLIDE_REQUIRED_TERMS"))
        self.assertTrue(hasattr(validator, "FIRST_PAGE_ML_BAN"))
        for term in (
            "AgentFit",
            "AgentTeams",
            "OpsPilot",
            "ProjectCase",
            "TaskSample",
            "Episode",
            EVALUATION_IDENTITY,
            "db_pool_exhausted",
            "slow_sql_degradation",
            "AgentSolutionPackage",
            "设计契约，非运行证据",
            "requires_runtime_trial",
            "同一搜索空间",
            "七层 ML 映射",
        ):
            with self.subTest(term=term):
                self.assertIn(term, validator.REQUIRED_TERMS)
        self.assertEqual(
            "完整方案空间：工具、Skill、MCP、Memory、模型、Agent 拓扑与 Human 边界都是变量。",
            validator.EXPECTED_PAGE_TITLES[4],
        )
        self.assertEqual(
            "五阶段闭环：定义案例与验收、构建最小候选、运行并测量、分析并调整、验证并停止。",
            validator.EXPECTED_PAGE_TITLES[5],
        )
        self.assertEqual(
            "证据账本：OpsPilot 为官方锚点，retail/airline 仅探索性 Demo。",
            validator.EXPECTED_PAGE_TITLES[10],
        )
        for term in (
            "完整方案空间",
            "五阶段闭环",
            validator.SLIDE_11_EVIDENCE_STATEMENT,
        ):
            with self.subTest(term=term):
                self.assertTrue(
                    any(
                        term in terms
                        for terms in validator.SLIDE_REQUIRED_TERMS.values()
                    )
                )

    def test_validator_forbids_overclaims(self) -> None:
        validator = _load_validator()
        for term in (
            "Agent 方案训练系统",
            "AutoML for Agents",
            "语义反向传播",
            "已选定 C2",
            "C2 胜出",
            "已跑通最小闭环",
            "ROI 提升",
            "准确率 9",
            "Meta-learning 已验证",
            "proxy 分数是官方结果",
            "代理分数是官方结果",
            "proxy score is official",
            "exploratory proxy scores are official results",
        ):
            with self.subTest(term=term):
                self.assertIn(term, validator.FORBIDDEN_TERMS)
        for term in (
            "AutoML",
            "semantic gradient",
            "语义梯度",
            "backpropagation",
            "反向传播",
            "official benchmark accuracy",
            "官方基准准确率",
            "正式 Candidate 执行完成",
            "formal Candidate execution completed",
            "AgentTeams 端到端集成完成",
            "AgentTeams end-to-end integration complete",
        ):
            with self.subTest(term=term):
                self.assertIn(term, validator.FORBIDDEN_TERMS)

    def test_validator_rejects_broadened_claim_mutations_in_pptx_and_pdf(self) -> None:
        validator = _load_validator()
        mutation_terms = (
            "AutoML",
            "semantic gradient",
            "backpropagation",
            "official benchmark accuracy",
            "formal Candidate execution completed",
            "AgentTeams end-to-end integration complete",
        )
        meta_learning_claims = (
            "Meta-learning is verified",
            "verified Meta-learning",
            "已验证 Meta-learning",
            "Meta learning is verified",
            "Meta-learning has been verified",
            "verified Meta learning",
            "VERIFIED META LEARNING",
        )
        payload = "\n".join(mutation_terms + meta_learning_claims)
        with tempfile.TemporaryDirectory(prefix="submission-overclaim-") as temp_dir:
            temp = Path(temp_dir)

            pptx_path = temp / "overclaim.pptx"
            presentation = Presentation(PPTX)
            textbox = presentation.slides[16].shapes.add_textbox(
                Inches(0.1), Inches(0.1), Inches(10), Inches(1.5)
            )
            textbox.text = payload
            presentation.save(pptx_path)
            pptx_errors = validator.validate(pptx_path)

            overlay_path = temp / "overlay.pdf"
            canvas = Canvas(str(overlay_path), pagesize=(1280, 720))
            canvas.setFont("Helvetica", 10)
            for line_number, line in enumerate(mutation_terms + meta_learning_claims):
                canvas.drawString(20, 700 - line_number * 16, line)
            canvas.save()

            pdf_path = temp / "overclaim.pdf"
            overlay_reader = PdfReader(overlay_path)
            writer = PdfWriter(clone_from=str(PDF))
            writer.pages[0].merge_page(overlay_reader.pages[0])
            with pdf_path.open("wb") as stream:
                writer.write(stream)
            pdf_errors = validator.validate(PPTX, pdf_path)

        for term in mutation_terms:
            with self.subTest(artifact="PPTX", term=term):
                self.assertIn(f"PPTX contains forbidden term: {term}", pptx_errors)
            with self.subTest(artifact="PDF", term=term):
                self.assertIn(f"PDF contains forbidden term: {term}", pdf_errors)
        for artifact, errors in (("PPTX", pptx_errors), ("PDF", pdf_errors)):
            with self.subTest(artifact=artifact, claim="verified Meta-learning"):
                self.assertEqual(
                    1,
                    errors.count(
                        f"{artifact} contains forbidden declaration pattern: "
                        "verified Meta-learning claim"
                    ),
                    errors,
                )

    def test_validator_rejects_extra_identity_and_skill_mutations(self) -> None:
        validator = _load_validator()
        identity_numbered_mutation = "6 / Coordinator Agent"
        identity_unnumbered_mutation = "Coordinator Agent"
        skill_numbered_mutation = "10 / 工具编排"
        skill_unnumbered_mutation = "工具编排"
        with tempfile.TemporaryDirectory(prefix="submission-entry-count-") as temp_dir:
            temp = Path(temp_dir)

            def mutate_pptx(
                path: Path, slide_number: int, anchor: str, mutation: str
            ) -> list[str]:
                presentation = Presentation(PPTX)
                matches = [
                    shape
                    for shape in validator._iter_shapes(
                        presentation.slides[slide_number - 1].shapes
                    )
                    if getattr(shape, "has_text_frame", False)
                    and anchor in shape.text
                ]
                self.assertEqual(
                    1,
                    len(matches),
                    f"frozen list anchor must identify one text shape: {anchor!r}",
                )
                matches[0].text = f"{matches[0].text}\n{mutation}"
                presentation.save(path)
                return validator.validate(path)

            identity_numbered_errors = mutate_pptx(
                temp / "extra-identity-numbered.pptx",
                7,
                "01 交付官",
                identity_numbered_mutation,
            )
            identity_unnumbered_errors = mutate_pptx(
                temp / "extra-identity-unnumbered.pptx",
                7,
                "01 交付官",
                identity_unnumbered_mutation,
            )
            skill_numbered_errors = mutate_pptx(
                temp / "extra-skill-numbered.pptx",
                15,
                "1 任务编译",
                skill_numbered_mutation,
            )
            skill_unnumbered_errors = mutate_pptx(
                temp / "extra-skill-unnumbered.pptx",
                15,
                "1 任务编译",
                skill_unnumbered_mutation,
            )

            def mutate_standalone_pptx(
                path: Path, slide_number: int, anchor: str, mutation: str
            ) -> list[str]:
                presentation = Presentation(PPTX)
                slide = presentation.slides[slide_number - 1]
                matches = [
                    shape
                    for shape in validator._iter_shapes(slide.shapes)
                    if getattr(shape, "has_text_frame", False)
                    and anchor in shape.text
                ]
                self.assertEqual(1, len(matches), anchor)
                anchor_shape = matches[0]
                textbox = slide.shapes.add_textbox(
                    anchor_shape.left,
                    anchor_shape.top,
                    anchor_shape.width,
                    anchor_shape.height,
                )
                textbox.text = mutation
                presentation.save(path)
                return validator.validate(path)

            identity_standalone_errors = mutate_standalone_pptx(
                temp / "extra-identity-standalone.pptx",
                7,
                "01 交付官",
                identity_unnumbered_mutation,
            )
            skill_standalone_errors = mutate_standalone_pptx(
                temp / "extra-skill-standalone.pptx",
                15,
                "1 任务编译",
                skill_unnumbered_mutation,
            )

            def pdf_anchor_position(
                page_number: int, number: int, label: str
            ) -> tuple[float, float, float]:
                page = PdfReader(PDF).pages[page_number - 1]
                spans: list[tuple[float, float, str, float]] = []

                def visitor(
                    text: str,
                    _cm: object,
                    tm: object,
                    _font: object,
                    size: float,
                ) -> None:
                    if text.strip():
                        spans.append((float(tm[4]), float(tm[5]), text, float(size)))

                page.extract_text(visitor_text=visitor)
                anchor_pattern = rf"^\s*0?{number}(?=\s|/|$)"
                matches = [
                    span
                    for span in spans
                    if re.match(anchor_pattern, span[2])
                    and any(
                        abs(span[1] - other[1]) <= 1
                        and abs(span[0] - other[0]) <= 120
                        and label[:1] in other[2]
                        for other in spans
                    )
                ]
                self.assertEqual(
                    1,
                    len(matches),
                    f"frozen PDF list anchor must identify one text span: {number} {label}",
                )
                return matches[0][0], matches[0][1], matches[0][3]

            pdfmetrics.registerFont(
                TTFont(
                    "DroidSansFallback",
                    "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
                )
            )

            def mutate_pdf(
                path: Path,
                page_number: int,
                number: int,
                label: str,
                mutation: str,
                unnumbered: bool = False,
            ) -> list[str]:
                base_reader = PdfReader(PDF)
                page = base_reader.pages[page_number - 1]
                x, y, size = pdf_anchor_position(page_number, number, label)
                overlay_path = temp / f"{path.stem}-overlay.pdf"
                canvas = Canvas(
                    str(overlay_path),
                    pagesize=(float(page.mediabox.width), float(page.mediabox.height)),
                )
                canvas.setFont("DroidSansFallback", size)
                for overlay_page in range(len(base_reader.pages)):
                    if overlay_page == page_number - 1:
                        draw_y = y - size * 1.5 if unnumbered else y
                        canvas.drawString(x, draw_y, mutation)
                    canvas.showPage()
                canvas.save()
                overlay_reader = PdfReader(overlay_path)
                writer = PdfWriter(clone_from=str(PDF))
                writer.pages[page_number - 1].merge_page(
                    overlay_reader.pages[page_number - 1]
                )
                with path.open("wb") as stream:
                    writer.write(stream)
                return validator.validate(PPTX, path)

            identity_numbered_pdf_errors = mutate_pdf(
                temp / "extra-identity-numbered.pdf",
                7,
                1,
                "交付官",
                identity_numbered_mutation,
            )
            identity_unnumbered_pdf_errors = mutate_pdf(
                temp / "extra-identity-unnumbered.pdf",
                7,
                1,
                "交付官",
                identity_unnumbered_mutation,
                True,
            )
            skill_numbered_pdf_errors = mutate_pdf(
                temp / "extra-skill-numbered.pdf",
                15,
                1,
                "任务编译",
                skill_numbered_mutation,
            )
            skill_unnumbered_pdf_errors = mutate_pdf(
                temp / "extra-skill-unnumbered.pdf",
                15,
                1,
                "任务编译",
                skill_unnumbered_mutation,
                True,
            )

        for errors in (
            identity_numbered_errors,
            identity_numbered_pdf_errors,
        ):
            with self.subTest(artifact_errors=errors):
                self.assertTrue(
                    any("exactly five fixed Meta Agent entries" in error for error in errors),
                    errors,
                )
        for errors in (
            identity_unnumbered_errors,
            skill_unnumbered_errors,
            identity_standalone_errors,
            skill_standalone_errors,
            identity_unnumbered_pdf_errors,
            skill_unnumbered_pdf_errors,
        ):
            with self.subTest(artifact_errors=errors):
                self.assertTrue(
                    any("unrecognized extra identity/Skill entry" in error for error in errors),
                    errors,
                )
        for errors in (skill_numbered_errors, skill_numbered_pdf_errors):
            with self.subTest(artifact_errors=errors):
                self.assertTrue(
                    any("exactly seven fixed Skill entries" in error for error in errors),
                    errors,
                )

    def test_validator_scopes_unexpected_entry_checks_to_list_regions(self) -> None:
        validator = _load_validator()
        skill_page = (SLIDES_DIR / "15-a3-skills.html").read_text(
            encoding="utf-8"
        )
        page_text = re.sub(r"<[^>]+>", "", skill_page)
        page_text += "\nadditional Skill is a normal explanatory phrase."
        errors = validator._numbered_contract_errors("HTML page", 15, page_text)
        self.assertNotIn(
            "HTML page contains an unrecognized extra identity/Skill entry",
            errors,
        )

        with tempfile.TemporaryDirectory(prefix="submission-list-copy-") as temp_dir:
            path = Path(temp_dir) / "rewritten-surrounding-copy.pptx"
            presentation = Presentation(PPTX)
            replacements = {
                14: {
                    "完整 Identity 契约": "Identity 字段说明",
                    "候选冻结后": "Supporting Agent is ordinary explanatory copy. 候选冻结后",
                },
                15: {
                    "SEVEN SKILLS": "CORE METHODS",
                    "SKILL · 可复用方法": "METHOD · 可复用能力",
                    "完整字段见": "additional Skill is ordinary explanatory copy. 完整字段见",
                },
            }
            for page_number, page_replacements in replacements.items():
                for shape in validator._iter_shapes(
                    presentation.slides[page_number - 1].shapes
                ):
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    for old, new in page_replacements.items():
                        if old in shape.text:
                            shape.text = shape.text.replace(old, new)
            presentation.save(path)
            pptx_errors = validator.validate(path)

        self.assertFalse(
            any(
                "unrecognized extra identity/Skill entry" in error
                for error in pptx_errors
            ),
            pptx_errors,
        )

    def test_validator_rejects_raster_media_and_transitions(self) -> None:
        validator = _load_validator()
        with tempfile.TemporaryDirectory(prefix="submission-native-") as temp_dir:
            temp = Path(temp_dir)
            pptx_path = temp / "non-native.pptx"
            png_path = temp / "pixel.png"
            png_path.write_bytes(
                base64.b64decode(
                    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
                )
            )
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(
                str(png_path),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
            presentation.save(pptx_path)
            staged = temp / "staged.pptx"
            with zipfile.ZipFile(pptx_path) as source, zipfile.ZipFile(
                staged, "w"
            ) as target:
                for item in source.infolist():
                    payload = source.read(item.filename)
                    if item.filename == "ppt/slides/slide1.xml":
                        payload = payload.replace(
                            b"</p:sld>", b"<p:transition/></p:sld>"
                        )
                    target.writestr(item, payload)
                target.writestr("ppt/media/demo.mp4", b"design-test-media")
            staged.replace(pptx_path)

            errors = validator.validate(pptx_path)

        self.assertIn("PPTX slide 1 contains a picture/raster shape", errors)
        self.assertIn("PPTX package contains embedded media: ppt/media/demo.mp4", errors)
        self.assertIn("PPTX slide 1 contains a transition", errors)

    def test_validator_flags_ml_term_on_first_four_slides(self) -> None:
        validator = _load_validator()
        self.assertTrue(PPTX.is_file(), "PPTX must be built before running this test")
        with tempfile.TemporaryDirectory(prefix="submission-ml-ban-") as temp_dir:
            path = Path(temp_dir) / "ml-on-page1.pptx"
            presentation = Presentation(PPTX)
            textbox = presentation.slides[0].shapes.add_textbox(
                Inches(0.1), Inches(0.1), Inches(10), Inches(0.4)
            )
            textbox.text = "这里用 Meta-learning 选 NAS 架构"
            presentation.save(path)
            errors = validator.validate(path)
        self.assertTrue(
            any("ML/NAS" in error and "slide 1" in error for error in errors),
            f"expected an ML/NAS ban error on slide 1, got: {errors}",
        )

    def test_validator_flags_missing_page_title(self) -> None:
        validator = _load_validator()
        self.assertTrue(PPTX.is_file())
        title = validator.EXPECTED_PAGE_TITLES[3]
        with tempfile.TemporaryDirectory(prefix="submission-title-") as temp_dir:
            path = Path(temp_dir) / "titled.pptx"
            presentation = Presentation(PPTX)
            for shape in presentation.slides[3].shapes:
                if getattr(shape, "has_text_frame", False) and title in shape.text:
                    shape.text = shape.text.replace(title, "占位标题")
            presentation.save(path)
            errors = validator.validate(path)
        self.assertIn(
            f"PPTX slide 4 is missing expected title: {title}",
            errors,
        )

    def test_built_pptx_and_pdf_validate_clean(self) -> None:
        validator = _load_validator()
        self.assertTrue(PPTX.is_file(), "PPTX must be built")
        self.assertTrue(PDF.is_file(), "PDF must be built")
        errors = validator.validate(PPTX, PDF)
        self.assertEqual([], errors, f"validation errors: {errors}")

    def test_validator_rejects_seventeen_blank_pdf_pages(self) -> None:
        validator = _load_validator()
        self.assertTrue(PPTX.is_file())
        with tempfile.TemporaryDirectory(prefix="submission-pdf-") as temp_dir:
            pdf_path = Path(temp_dir) / "blank-17.pdf"
            writer = PdfWriter()
            for _ in range(17):
                writer.add_blank_page(width=1280, height=720)
            with pdf_path.open("wb") as stream:
                writer.write(stream)
            errors = validator.validate(PPTX, pdf_path)
        title = validator.EXPECTED_PAGE_TITLES[0]
        self.assertIn(
            f"PDF page 1 is missing expected title: {title}",
            errors,
        )

    def test_validator_rejects_pdf_overclaim_and_missing_pptx_text(self) -> None:
        validator = _load_validator()
        with tempfile.TemporaryDirectory(prefix="submission-pdf-parity-") as temp_dir:
            pdf_path = Path(temp_dir) / "divergent.pdf"
            canvas = Canvas(str(pdf_path), pagesize=(1280, 720))
            for index, title in enumerate(validator.EXPECTED_PAGE_TITLES, start=1):
                canvas.setFont("Helvetica", 8)
                canvas.drawString(20, 700, f"page {index}")
                if index == 1:
                    canvas.drawString(20, 680, "ImageNet")
                canvas.showPage()
            canvas.save()

            errors = validator.validate(PPTX, pdf_path)

        self.assertIn("PDF contains forbidden term: ImageNet", errors)
        self.assertTrue(
            any(
                error.startswith("PDF page 1 is missing PPTX text:")
                for error in errors
            ),
            errors,
        )


if __name__ == "__main__":
    unittest.main()
