#!/usr/bin/env python3
"""Validate the AgentFit scoreline-v2 alternative deck deliverables."""

from __future__ import annotations

import argparse
import re
import sys
import zipfile
from xml.etree import ElementTree
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader


EXPECTED_SLIDES = 17

EXPECTED_PAGE_TITLES = (
    "选对 Agent 方案，不该靠猜。",
    "OpsPilot 官方 baseline 是首个 ProjectCase 参考。",
    "AgentFit 把业务难题编译成 ProjectCase。",
    "一次事故定义 TaskSample、输入输出与验收。",
    "五元元团队完成方案设计闭环。",
    "C0 / C1 / C2 / C3 是同一搜索空间的候选。",
    "五元团队映射到 AgentTeams 的 Worker / Team / Room / Human。",
    "公平试验、Episode、Trace 与证据契约尚待真实运行。",
    "七个 Skill、HTTP/MCP 契约、共享状态与风险门禁支撑闭环。",
    "交付 AgentSolutionPackage 与五种合法结果。",
    "证据账本：baseline 已代码级审计，候选对照仍待运行。",
    "Fit：不是更多 Agent，而是有证据选对方案。",
    "候选表示与搜索规则：不让 ML/NAS 抢主线。",
    "五个 Agent Identity：判断权、状态边界与责任产物。",
    "七个 Skill、HTTP/MCP 等价工具与上下文机制。",
    "Human 门禁、风险、异常与回滚。",
    "开放、依赖、许可证、baseline 引用与未实现边界。",
)

SLIDE_REQUIRED_TERMS = {
    2: ("OpsPilot", "ProjectCase", "baseline"),
    4: ("TaskSample", "db_pool_exhausted", "slow_sql_degradation"),
    7: ("Worker", "Team", "Room", "Human"),
    8: ("Episode", "Trace", "CandidateVersion × SampleVersion × RunIndex"),
    10: ("AgentSolutionPackage",),
    11: ("代码级审计", "仍待运行"),
    13: ("七层 ML 映射",),
    14: (
        "EngagementLead",
        "BusinessEngineer",
        "AgentArchitect",
        "ValidationEngineer",
        "GovernanceAuditor",
    ),
}

FIRST_PAGE_ML_BAN = (
    "Architecture Search",
    "NAS",
    "Meta-learning",
    "Meta learning",
    "七层",
    "L1",
    "L2",
    "L3",
    "SVD",
    "argmin",
    "Pareto",
    "贝叶斯",
    "inner loop",
    "outer loop",
)

REQUIRED_TERMS = (
    "AgentFit",
    "AgentTeams",
    "OpsPilot",
    "ProjectCase",
    "db_pool_exhausted",
    "slow_sql_degradation",
    "TaskSample",
    "Episode",
    "CandidateVersion × SampleVersion × RunIndex",
    "EngagementLead",
    "BusinessEngineer",
    "AgentArchitect",
    "ValidationEngineer",
    "GovernanceAuditor",
    "Agentless",
    "C0",
    "C1",
    "C2",
    "C3",
    "Human",
    "Skill",
    "MCP",
    "HTTP",
    "上下文",
    "验证",
    "安全",
    "开放",
    "未实现",
    "AgentSolutionPackage",
    "设计契约，非运行证据",
    "requires_runtime_trial",
    "Worker",
    "Team",
    "Room",
    "七层 ML 映射",
    "同一搜索空间",
)

FORBIDDEN_TERMS = (
    "已选定 C2",
    "已选定 C1",
    "C2 胜出",
    "已跑通最小闭环",
    "ROI 提升",
    "准确率 9",
    "ImageNet",
    "90%+",
    "Meta-learning 已验证",
    "已开放",
)

PRESENTATIONML_NAMESPACE = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _shape_text(shape: object) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        texts.append(shape.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                texts.append(cell.text)
    return texts


def _iter_shapes(shapes: object):
    for shape in shapes:
        yield shape
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_shapes(nested_shapes)


def _slide_text(slide: object) -> str:
    return "\n".join(
        text
        for shape in _iter_shapes(slide.shapes)
        for text in _shape_text(shape)
        if text.strip()
    )


def _normalized(text: str) -> str:
    return re.sub(r"\s+", "", text)


def _has_visible_native_content(
    slide: object, slide_width: int, slide_height: int
) -> bool:
    for shape in _iter_shapes(slide.shapes):
        if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.MEDIA):
            continue
        left = getattr(shape, "left", 0)
        top = getattr(shape, "top", 0)
        width = getattr(shape, "width", 0)
        height = getattr(shape, "height", 0)
        if width <= 0 or height <= 0:
            continue
        if left >= slide_width or top >= slide_height:
            continue
        if left + width <= 0 or top + height <= 0:
            continue
        if any(_normalized(text) for text in _shape_text(shape)):
            return True
    return False


def _package_editability_errors(pptx_path: Path) -> list[str]:
    errors: list[str] = []
    try:
        with zipfile.ZipFile(pptx_path) as package:
            media = sorted(
                name
                for name in package.namelist()
                if name.startswith("ppt/media/") and not name.endswith("/")
            )
            for name in media:
                errors.append(f"PPTX package contains embedded media: {name}")
            for slide_number in range(1, EXPECTED_SLIDES + 1):
                member = f"ppt/slides/slide{slide_number}.xml"
                if member not in package.namelist():
                    continue
                root = ElementTree.fromstring(package.read(member))
                transition = root.find(
                    f"{{{PRESENTATIONML_NAMESPACE}}}transition"
                )
                if transition is not None:
                    errors.append(f"PPTX slide {slide_number} contains a transition")
    except (OSError, zipfile.BadZipFile, ElementTree.ParseError) as exc:
        errors.append(f"Unable to inspect PPTX package {pptx_path}: {exc}")
    return errors


def validate(pptx_path: Path, pdf_path: Path | None = None) -> list[str]:
    """Return human-readable validation errors; an empty list means success."""
    errors: list[str] = []
    if not pptx_path.is_file():
        return [f"PPTX file does not exist: {pptx_path}"]

    try:
        presentation = Presentation(pptx_path)
    except Exception as exc:  # pragma: no cover
        return [f"Unable to read PPTX {pptx_path}: {exc}"]

    slide_count = len(presentation.slides)
    if slide_count != EXPECTED_SLIDES:
        errors.append(
            f"PPTX must contain exactly {EXPECTED_SLIDES} slides; found {slide_count}"
        )

    slide_texts = [_slide_text(slide) for slide in presentation.slides]
    for index, (slide, text) in enumerate(
        zip(presentation.slides, slide_texts), start=1
    ):
        if not text.strip():
            errors.append(f"PPTX slide {index} is empty")
        if not _has_visible_native_content(
            slide, presentation.slide_width, presentation.slide_height
        ):
            errors.append(
                f"PPTX slide {index} has no visible native text/content shape"
            )
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                errors.append(f"PPTX slide {index} contains a picture/raster shape")
            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                errors.append(f"PPTX slide {index} contains a media shape")

        normalized_text = _normalized(text)
        if index <= len(EXPECTED_PAGE_TITLES):
            title = EXPECTED_PAGE_TITLES[index - 1]
            if _normalized(title) not in normalized_text:
                errors.append(
                    f"PPTX slide {index} is missing expected title: {title}"
                )
        for term in SLIDE_REQUIRED_TERMS.get(index, ()):
            if _normalized(term) not in normalized_text:
                errors.append(
                    f"PPTX slide {index} is missing required term: {term}"
                )
        if index <= 4:
            for term in FIRST_PAGE_ML_BAN:
                if _normalized(term) in normalized_text:
                    errors.append(
                        f"PPTX slide {index} uses an ML/NAS term forbidden in the first four pages: {term}"
                    )

    errors.extend(_package_editability_errors(pptx_path))

    deck_text = "\n".join(slide_texts)
    for term in REQUIRED_TERMS:
        if term not in deck_text:
            errors.append(f"PPTX is missing required term: {term}")
    for term in FORBIDDEN_TERMS:
        if term in deck_text:
            errors.append(f"PPTX contains forbidden term: {term}")

    if pdf_path is not None:
        if not pdf_path.is_file():
            errors.append(f"PDF file does not exist: {pdf_path}")
        else:
            try:
                pdf_reader = PdfReader(pdf_path)
                pdf_pages = len(pdf_reader.pages)
            except Exception as exc:  # pragma: no cover
                errors.append(f"Unable to read PDF {pdf_path}: {exc}")
            else:
                if pdf_pages != EXPECTED_SLIDES:
                    errors.append(
                        f"PDF must contain exactly {EXPECTED_SLIDES} pages; found {pdf_pages}"
                    )
                for index, page in enumerate(pdf_reader.pages, start=1):
                    try:
                        page_text = page.extract_text() or ""
                    except Exception as exc:  # pragma: no cover
                        errors.append(f"Unable to extract PDF page {index} text: {exc}")
                        continue
                    normalized_page = _normalized(page_text)
                    if index <= len(EXPECTED_PAGE_TITLES):
                        title = EXPECTED_PAGE_TITLES[index - 1]
                        if _normalized(title) not in normalized_page:
                            errors.append(
                                f"PDF page {index} is missing expected title: {title}"
                            )
                    for term in SLIDE_REQUIRED_TERMS.get(index, ()):
                        if _normalized(term) not in normalized_page:
                            errors.append(
                                f"PDF page {index} is missing required term: {term}"
                            )
                    if index <= 4:
                        for term in FIRST_PAGE_ML_BAN:
                            if _normalized(term) in normalized_page:
                                errors.append(
                                    f"PDF page {index} uses an ML/NAS term forbidden in the first four pages: {term}"
                                )

    return errors


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Validate AgentFit scoreline-v2 PPTX and optional PDF."
    )
    parser.add_argument("pptx", type=Path, help="Path to the PPTX file")
    parser.add_argument("pdf", type=Path, nargs="?", help="Optional PDF export")
    args = parser.parse_args()

    errors = validate(args.pptx, args.pdf)
    if errors:
        print("VALIDATION_FAILED")
        for error in errors:
            print(f"- {error}")
        return 1

    presentation = Presentation(args.pptx)
    print(f"pptx_pages={len(presentation.slides)}")
    if args.pdf is not None:
        print(f"pdf_pages={len(PdfReader(args.pdf).pages)}")
    print("content_checks=PASS")
    print("native_editability_checks=PASS")
    if args.pdf is not None:
        print("pdf_page_text_checks=PASS")
    return 0


if __name__ == "__main__":
    sys.exit(main())
