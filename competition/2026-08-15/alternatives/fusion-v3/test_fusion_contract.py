#!/usr/bin/env python3
"""Contract tests for the AgentFit fusion-v3 alternative deck.

These tests run strictly against competition/2026-08-15/alternatives/fusion-v3/
and never touch the frozen submission/, scoreline-v2/ or canonical docs.
"""

from __future__ import annotations

import base64
import importlib.util
import re
import subprocess
import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from reportlab.pdfgen.canvas import Canvas


ROOT = Path(__file__).resolve().parent
SLIDES_DIR = ROOT / "slides"
INTRODUCTION = ROOT / "work-introduction.md"
OUTLINE = ROOT / "ppt-outline.md"
VALIDATOR = ROOT / "validate_presentation.py"
BUILDER = ROOT / "build_presentation.py"
PPTX = ROOT / "agentfit-fusion-v3.pptx"
PDF = ROOT / "agentfit-fusion-v3.pdf"

EVALUATION_IDENTITY = "CandidateVersion × SampleVersion × RunIndex"


def _load_validator():
    spec = importlib.util.spec_from_file_location("validate_presentation", VALIDATOR)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"unable to import {VALIDATOR}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _introduction_body(text: str) -> str:
    match = re.search(
        r"^## 500 字以内作品简介\s*$\n(?P<body>.*?)(?=^##\s|\Z)",
        text,
        flags=re.MULTILINE | re.DOTALL,
    )
    if not match:
        raise AssertionError("missing '500 字以内作品简介' section")
    return match.group("body").strip()


def _pdf_layout_lines(page: int) -> list[str]:
    completed = subprocess.run(
        [
            "pdftotext",
            "-f",
            str(page),
            "-l",
            str(page),
            "-layout",
            str(PDF),
            "-",
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    return [re.sub(r"\s+", " ", line).strip() for line in completed.stdout.splitlines()]


def _forward_arrows_between(slide: object, source: object, target: object) -> list[object]:
    return [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text.strip() == "→"
        and source.left + source.width <= shape.left
        and shape.left + shape.width <= target.left
    ]


class FusionV3ContractTest(unittest.TestCase):
    def test_directory_is_isolated_from_submission_and_scoreline(self) -> None:
        self.assertTrue(SLIDES_DIR.is_dir())
        self.assertTrue(VALIDATOR.is_file())
        self.assertTrue(BUILDER.is_file())
        self.assertTrue(INTRODUCTION.is_file())
        self.assertTrue(OUTLINE.is_file())
        root = str(ROOT)
        for path in ROOT.rglob("*"):
            self.assertTrue(
                str(path).startswith(root),
                f"stray file outside fusion-v3: {path}",
            )

    def test_seventeen_html_slides_with_stable_names(self) -> None:
        slides = sorted(SLIDES_DIR.glob("[0-9][0-9]-*.html"))
        self.assertEqual(17, len(slides))
        self.assertEqual(
            [f"{index:02d}" for index in range(1, 18)],
            [path.name[:2] for path in slides],
        )

    def test_html_sources_use_pptx_safe_solid_colors(self) -> None:
        for path in sorted(SLIDES_DIR.glob("*.html")):
            source = path.read_text(encoding="utf-8")
            with self.subTest(slide=path.name):
                self.assertNotRegex(source, r"\b(?:rgba|hsla)\s*\(")

    def test_final_convergence_copy_and_diagram_contract(self) -> None:
        cover = (SLIDES_DIR / "01-cover.html").read_text(encoding="utf-8")
        selection = (SLIDES_DIR / "06-selection-rule.html").read_text(
            encoding="utf-8"
        )
        search = (SLIDES_DIR / "13-a1-search.html").read_text(encoding="utf-8")
        risk = (SLIDES_DIR / "16-a4-risk.html").read_text(encoding="utf-8")

        self.assertIn("C2 · N 业务 Agent", cover)

        self.assertIn("候选搜索顺序 · 设计契约（非运行结果）", selection)
        for label in (
            "C0 · Agentless · 待真实试验",
            "C1 · 单 Agent · 待真实试验",
            "C2 · 多 Agent · 待真实试验",
            "C3 · Human 混合 · 待真实试验",
        ):
            with self.subTest(label=label):
                self.assertIn(label, selection)
        for misleading_mark in ("✓", "⚠", "验收阈值", "非真实数据"):
            with self.subTest(misleading_mark=misleading_mark):
                self.assertNotIn(misleading_mark, selection)

        self.assertIn("生成 → 定位 · 最多 3 次 · 未通过退出", search)
        self.assertIn("固定 G / Π / ρ；只更新 θ 与分区内状态", search)
        self.assertIn("更新 G / Π / ρ；比较 validation samples", search)

        self.assertIn("六类异常进入同一 Trace 语义：", risk)
        self.assertIn("对应恢复：", risk)

    def test_final_convergence_pdf_labels_do_not_wrap(self) -> None:
        expected_by_page = {
            1: ("C2 · N 业务 Agent",),
            6: (
                "候选搜索顺序 · 设计契约（非运行结果）",
                "C0 · Agentless · 待真实试验",
                "C1 · 单 Agent · 待真实试验",
                "C2 · 多 Agent · 待真实试验",
                "C3 · Human 混合 · 待真实试验",
            ),
            16: (
                "异常与恢复 · 同一 Trace 语义",
                "六类异常进入同一 Trace 语义：",
            ),
        }
        for page, labels in expected_by_page.items():
            lines = _pdf_layout_lines(page)
            for label in labels:
                with self.subTest(page=page, label=label):
                    self.assertTrue(any(label in line for line in lines), lines)

        page_13 = _pdf_layout_lines(13)
        self.assertFalse(
            any(line == "据。" for line in page_13),
            f"page 13 contains an orphaned final character: {page_13}",
        )

    def test_final_candidate_graph_has_no_text_overlap(self) -> None:
        presentation = Presentation(PPTX)
        slide = presentation.slides[12]
        by_text = {
            shape.text.strip(): shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        }
        arrow = by_text["←"]
        exit_condition = by_text["生成 → 定位 · 最多 3 次 · 未通过退出"]
        self.assertLessEqual(arrow.top + arrow.height, exit_condition.top)

    def test_final_candidate_graph_and_blocks_geometry(self) -> None:
        presentation = Presentation(PPTX)

        selection_slide = presentation.slides[5]
        candidate_texts = (
            "C0 · Agentless · 待真实试验",
            "C1 · 单 Agent · 待真实试验",
            "C2 · 多 Agent · 待真实试验",
            "C3 · Human 混合 · 待真实试验",
        )
        candidate_text_shapes = [
            next(shape for shape in selection_slide.shapes if shape.text == text)
            for text in candidate_texts
        ]
        candidate_blocks = [
            next(
                shape
                for shape in selection_slide.shapes
                if not getattr(shape, "text", "").strip()
                and shape.left <= text_shape.left
                and text_shape.left + text_shape.width <= shape.left + shape.width
                and shape.top <= text_shape.top
                and text_shape.top + text_shape.height <= shape.top + shape.height
                and shape.width < Inches(3)
            )
            for text_shape in candidate_text_shapes
        ]
        self.assertEqual(1, len({shape.width for shape in candidate_blocks}))
        self.assertEqual(1, len({shape.height for shape in candidate_blocks}))
        self.assertEqual(1, len({shape.top for shape in candidate_blocks}))

        search_slide = presentation.slides[12]
        by_text = {
            shape.text.strip(): shape
            for shape in search_slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        }
        forward_edges = (
            ("去重", "检索"),
            ("检索", "定位"),
            ("定位", "生成"),
            ("生成", "测试"),
            ("测试", "人工"),
        )
        for source, target in forward_edges:
            self.assertLess(by_text[source].left, by_text[target].left)
            arrows = _forward_arrows_between(
                search_slide, by_text[source], by_text[target]
            )
            self.assertEqual(
                1,
                len(arrows),
                f"expected exactly one forward arrow for {source} → {target}",
            )

            mutated = Presentation(PPTX)
            mutated_slide = mutated.slides[12]
            mutated_by_text = {
                shape.text.strip(): shape
                for shape in mutated_slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            }
            mutated_arrow = _forward_arrows_between(
                mutated_slide, mutated_by_text[source], mutated_by_text[target]
            )[0]
            mutated_arrow._element.getparent().remove(mutated_arrow._element)
            self.assertEqual(
                [],
                _forward_arrows_between(
                    mutated_slide, mutated_by_text[source], mutated_by_text[target]
                ),
                f"mutation proof failed for removed edge {source} → {target}",
            )

        locator = by_text["定位"]
        generator = by_text["生成"]

        return_arrow = by_text["←"]
        return_line = next(
            shape
            for shape in search_slide.shapes
            if not getattr(shape, "text", "").strip()
            and shape.height <= Inches(0.03)
            and locator.left < shape.left
            and shape.left < generator.left + generator.width
            and shape.top > locator.top + locator.height
        )
        self.assertLess(return_arrow.left, return_line.left)
        self.assertGreater(return_line.width, Inches(1.0))

    def test_introduction_under_five_hundred_non_whitespace_chars(self) -> None:
        body = _introduction_body(INTRODUCTION.read_text(encoding="utf-8"))
        count = len(re.sub(r"\s+", "", body))
        self.assertLessEqual(count, 500)
        for term in ("AgentFit", "AgentTeams", "OpsPilot", "方案建筑师"):
            self.assertIn(term, body)

    def test_first_four_slides_avoid_ml_and_nas_terminology(self) -> None:
        banned = (
            "Architecture Search",
            "NAS",
            "Meta-learning",
            "Meta learning",
            "七层",
            "L1",
            "L2",
            "L3",
            "SVD",
            "argmin",
            "Pareto",
            "贝叶斯",
            "inner loop",
            "outer loop",
        )
        for index in range(1, 5):
            matches = sorted(SLIDES_DIR.glob(f"{index:02d}-*.html"))
            self.assertEqual(1, len(matches), f"expected one slide for page {index}")
            text = matches[0].read_text(encoding="utf-8")
            for term in banned:
                with self.subTest(page=index, term=term):
                    self.assertNotIn(term, text)

    def test_cover_keeps_general_product_position_not_ops_product(self) -> None:
        cover = (SLIDES_DIR / "01-cover.html").read_text(encoding="utf-8")
        for term in ("方案建筑师", "AgentTeams"):
            self.assertIn(term, cover)
        for term in ("运维产品", "OpsPilot Zero 是 AgentFit"):
            self.assertNotIn(term, cover)

    def test_meta_agent_vs_business_agent_two_layers_are_explicit(self) -> None:
        meta_team = (SLIDES_DIR / "07-meta-team.html").read_text(encoding="utf-8")
        for term in (
            "EngagementLead",
            "BusinessEngineer",
            "AgentArchitect",
            "ValidationEngineer",
            "GovernanceAuditor",
            "元 Agent",
            "业务执行 Agent",
        ):
            self.assertIn(term, meta_team)

    def test_compact_pdf_labels_do_not_wrap_orphan_characters(self) -> None:
        page_7 = _pdf_layout_lines(7)
        for label in (
            "01 交付官",
            "02 业务架构师",
            "03 方案架构师",
            "04 验证工程师",
            "05 审计官",
        ):
            with self.subTest(page=7, label=label):
                self.assertTrue(any(label in line for line in page_7), page_7)

        page_9 = _pdf_layout_lines(9)
        for label in (
            "共享状态 · 上下文 4 选 2",
            "示例：读配置 · 列变更 · 写回滚",
        ):
            with self.subTest(page=9, label=label):
                self.assertTrue(any(label in line for line in page_9), page_9)

    def test_baseline_incidents_and_task_sample_terms_are_present(self) -> None:
        source = "\n".join(
            path.read_text(encoding="utf-8")
            for path in sorted(SLIDES_DIR.glob("*.html"))
        )
        for term in (
            "OpsPilot",
            "ProjectCase",
            "db_pool_exhausted",
            "slow_sql_degradation",
            "TaskSample",
            "Episode",
            EVALUATION_IDENTITY,
            "设计契约，非运行证据",
            "requires_runtime_trial",
            "同一搜索空间",
        ):
            with self.subTest(term=term):
                self.assertIn(term, source)

    def test_search_space_and_seven_layer_mapping_restored(self) -> None:
        source = "\n".join(
            path.read_text(encoding="utf-8")
            for path in sorted(SLIDES_DIR.glob("*.html"))
        )
        for term in (
            "七层 ML 映射",
            "G, Π, θ, ρ",
            "inner loop",
            "outer loop",
            "Meta-learning",
        ):
            with self.subTest(term=term):
                self.assertIn(term, source)

    def test_no_fabricated_runtime_evidence_in_html(self) -> None:
        source = "\n".join(
            path.read_text(encoding="utf-8")
            for path in sorted(SLIDES_DIR.glob("*.html"))
        )
        forbidden = (
            "已选定 C2",
            "已选定 C1",
            "已跑通最小闭环",
            "C2 胜出",
            "ROI 提升",
            "准确率 9",
            "ImageNet",
            "90%+",
            "已开放",
            "Meta-learning 已验证",
        )
        for term in forbidden:
            with self.subTest(term=term):
                self.assertNotIn(term, source)

    def test_validator_encodes_fusion_contract(self) -> None:
        validator = _load_validator()
        self.assertEqual(17, validator.EXPECTED_SLIDES)
        self.assertTrue(hasattr(validator, "EXPECTED_PAGE_TITLES"))
        self.assertEqual(17, len(validator.EXPECTED_PAGE_TITLES))
        self.assertTrue(hasattr(validator, "SLIDE_REQUIRED_TERMS"))
        self.assertTrue(hasattr(validator, "FIRST_PAGE_ML_BAN"))
        for term in (
            "AgentFit",
            "AgentTeams",
            "OpsPilot",
            "ProjectCase",
            "TaskSample",
            "Episode",
            EVALUATION_IDENTITY,
            "db_pool_exhausted",
            "slow_sql_degradation",
            "AgentSolutionPackage",
            "设计契约，非运行证据",
            "requires_runtime_trial",
            "同一搜索空间",
            "七层 ML 映射",
        ):
            with self.subTest(term=term):
                self.assertIn(term, validator.REQUIRED_TERMS)

    def test_validator_forbids_overclaims(self) -> None:
        validator = _load_validator()
        for term in (
            "已选定 C2",
            "C2 胜出",
            "已跑通最小闭环",
            "ROI 提升",
            "准确率 9",
            "Meta-learning 已验证",
        ):
            with self.subTest(term=term):
                self.assertIn(term, validator.FORBIDDEN_TERMS)

    def test_validator_rejects_raster_media_and_transitions(self) -> None:
        validator = _load_validator()
        with tempfile.TemporaryDirectory(prefix="fusion-native-") as temp_dir:
            temp = Path(temp_dir)
            pptx_path = temp / "non-native.pptx"
            png_path = temp / "pixel.png"
            png_path.write_bytes(
                base64.b64decode(
                    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
                )
            )
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_picture(
                str(png_path),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
            presentation.save(pptx_path)
            staged = temp / "staged.pptx"
            with zipfile.ZipFile(pptx_path) as source, zipfile.ZipFile(
                staged, "w"
            ) as target:
                for item in source.infolist():
                    payload = source.read(item.filename)
                    if item.filename == "ppt/slides/slide1.xml":
                        payload = payload.replace(
                            b"</p:sld>", b"<p:transition/></p:sld>"
                        )
                    target.writestr(item, payload)
                target.writestr("ppt/media/demo.mp4", b"design-test-media")
            staged.replace(pptx_path)

            errors = validator.validate(pptx_path)

        self.assertIn("PPTX slide 1 contains a picture/raster shape", errors)
        self.assertIn("PPTX package contains embedded media: ppt/media/demo.mp4", errors)
        self.assertIn("PPTX slide 1 contains a transition", errors)

    def test_validator_flags_ml_term_on_first_four_slides(self) -> None:
        validator = _load_validator()
        self.assertTrue(PPTX.is_file(), "PPTX must be built before running this test")
        with tempfile.TemporaryDirectory(prefix="fusion-ml-ban-") as temp_dir:
            path = Path(temp_dir) / "ml-on-page1.pptx"
            presentation = Presentation(PPTX)
            textbox = presentation.slides[0].shapes.add_textbox(
                Inches(0.1), Inches(0.1), Inches(10), Inches(0.4)
            )
            textbox.text = "这里用 Meta-learning 选 NAS 架构"
            presentation.save(path)
            errors = validator.validate(path)
        self.assertTrue(
            any("ML/NAS" in error and "slide 1" in error for error in errors),
            f"expected an ML/NAS ban error on slide 1, got: {errors}",
        )

    def test_validator_flags_missing_page_title(self) -> None:
        validator = _load_validator()
        self.assertTrue(PPTX.is_file())
        title = validator.EXPECTED_PAGE_TITLES[3]
        with tempfile.TemporaryDirectory(prefix="fusion-title-") as temp_dir:
            path = Path(temp_dir) / "titled.pptx"
            presentation = Presentation(PPTX)
            for shape in presentation.slides[3].shapes:
                if getattr(shape, "has_text_frame", False) and title in shape.text:
                    shape.text = shape.text.replace(title, "占位标题")
            presentation.save(path)
            errors = validator.validate(path)
        self.assertIn(
            f"PPTX slide 4 is missing expected title: {title}",
            errors,
        )

    def test_built_pptx_and_pdf_validate_clean(self) -> None:
        validator = _load_validator()
        self.assertTrue(PPTX.is_file(), "PPTX must be built")
        self.assertTrue(PDF.is_file(), "PDF must be built")
        errors = validator.validate(PPTX, PDF)
        self.assertEqual([], errors, f"validation errors: {errors}")

    def test_validator_rejects_seventeen_blank_pdf_pages(self) -> None:
        validator = _load_validator()
        self.assertTrue(PPTX.is_file())
        with tempfile.TemporaryDirectory(prefix="fusion-pdf-") as temp_dir:
            pdf_path = Path(temp_dir) / "blank-17.pdf"
            writer = PdfWriter()
            for _ in range(17):
                writer.add_blank_page(width=1280, height=720)
            with pdf_path.open("wb") as stream:
                writer.write(stream)
            errors = validator.validate(PPTX, pdf_path)
        title = validator.EXPECTED_PAGE_TITLES[0]
        self.assertIn(
            f"PDF page 1 is missing expected title: {title}",
            errors,
        )

    def test_validator_rejects_pdf_overclaim_and_missing_pptx_text(self) -> None:
        validator = _load_validator()
        with tempfile.TemporaryDirectory(prefix="fusion-pdf-parity-") as temp_dir:
            pdf_path = Path(temp_dir) / "divergent.pdf"
            canvas = Canvas(str(pdf_path), pagesize=(1280, 720))
            for index, title in enumerate(validator.EXPECTED_PAGE_TITLES, start=1):
                canvas.setFont("Helvetica", 8)
                canvas.drawString(20, 700, f"page {index}")
                if index == 1:
                    canvas.drawString(20, 680, "ImageNet")
                canvas.showPage()
            canvas.save()

            errors = validator.validate(PPTX, pdf_path)

        self.assertIn("PDF contains forbidden term: ImageNet", errors)
        self.assertTrue(
            any(
                error.startswith("PDF page 1 is missing PPTX text:")
                for error in errors
            ),
            errors,
        )


if __name__ == "__main__":
    unittest.main()
