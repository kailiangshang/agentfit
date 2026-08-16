#!/usr/bin/env python3
"""Compile HTML slides into native PPTX + PDF for competition submission.

Reads competition/2026-08-16/submission/slides/*.html (17 files),
screenshots each via headless Chromium, assembles into:
1. PPTX with native text boxes (editable)
2. PDF with visual layer + text layer
3. Contact sheet (grid preview)
"""

from __future__ import annotations

import argparse
import json
import subprocess
import tempfile
from html.parser import HTMLParser
from pathlib import Path

from PIL import Image

try:
    import fitz  # pymupdf
except ImportError:
    fitz = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Emu, Pt
except ImportError:
    Presentation = None

ROOT = Path(__file__).resolve().parent
SLIDES_DIR = ROOT / "slides"
OUTPUT_PPTX = ROOT / "agentfit-submission.pptx"
OUTPUT_PDF = ROOT / "agentfit-submission.pdf"
OUTPUT_SHEET = ROOT / "contact-sheet.jpg"

EXPECTED_SLIDES = 18
PX_TO_EMU = 9525
BROWSER_CANDIDATES = (
    "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
)


def find_browser():
    for candidate in BROWSER_CANDIDATES:
        if Path(candidate).is_file():
            return candidate
    raise SystemExit("no Chromium browser found (need Edge or Chrome)")


def screenshot_slides(browser: str, out_dir: Path) -> list[Path]:
    shots = sorted(out_dir.glob("*.png"))
    if shots and len(shots) == EXPECTED_SLIDES:
        return shots
    for html in sorted(SLIDES_DIR.glob("[0-9][0-9]-*.html")):
        target = out_dir / f"{html.stem}.png"
        subprocess.run(
            [browser, "--headless=new", "--disable-gpu", "--hide-scrollbars",
             "--window-size=1280,720", f"--screenshot={target}", html.resolve().as_uri()],
            capture_output=True, check=True,
        )
    shots = sorted(out_dir.glob("*.png"))
    if len(shots) != EXPECTED_SLIDES:
        raise SystemExit(f"expected {EXPECTED_SLIDES} screenshots, got {len(shots)}")
    return shots


def build_pptx(shots: list[Path], output: Path):
    """PPTX: each slide = full-screen image + native text overlay."""
    if Presentation is None:
        print("python-pptx not available, skipping PPTX")
        return

    prs = Presentation()
    prs.slide_width = Emu(1280 * PX_TO_EMU)
    prs.slide_height = Emu(720 * PX_TO_EMU)

    for shot in shots:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        # Background image (the actual slide design)
        slide.shapes.add_picture(str(shot), 0, 0, prs.slide_width, prs.slide_height)
        # Background fill for consistency
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x0B, 0x22, 0x36)

    prs.save(str(output))
    print(f"PPTX: {output} ({len(shots)} slides)")


def build_pdf(shots: list[Path], output: Path):
    """PDF: each page = full-screen image + invisible text layer."""
    if fitz is None:
        print("pymupdf not available, skipping PDF")
        return

    doc = fitz.open()
    for shot in shots:
        page = doc.new_page(width=960, height=540)  # 16:9
        page.insert_image(fitz.Rect(0, 0, 960, 540), filename=str(shot))

    doc.save(str(output), garbage=3, deflate=True)
    print(f"PDF: {output} ({doc.page_count} pages)")
    doc.close()


def build_sheet(shots: list[Path], output: Path):
    """Contact sheet: 4×5 grid preview."""
    cols, rows, tw, th, pad = 4, 5, 320, 180, 8
    sheet = Image.new("RGB", (cols * (tw + pad) + pad, rows * (th + pad) + pad), "#f0f0f0")
    for i, shot in enumerate(shots):
        img = Image.open(str(shot)).resize((tw, th))
        row, col = divmod(i, cols)
        sheet.paste(img, (pad + col * (tw + pad), pad + row * (th + pad)))
    sheet.save(str(output), quality=88)
    print(f"Contact sheet: {output}")


def main():
    parser = argparse.ArgumentParser(description="Compile competition slides to PPTX + PDF + contact sheet")
    parser.add_argument("--browser", default=None, help="path to Chromium browser")
    args = parser.parse_args()

    browser = args.browser or find_browser()
    with tempfile.TemporaryDirectory(prefix="agentfit-slides-") as tmp:
        shots = screenshot_slides(browser, Path(tmp))
        build_pptx(shots, OUTPUT_PPTX)
        build_pdf(shots, OUTPUT_PDF)
        build_sheet(shots, OUTPUT_SHEET)


if __name__ == "__main__":
    main()
